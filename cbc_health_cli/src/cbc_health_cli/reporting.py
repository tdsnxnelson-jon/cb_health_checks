from __future__ import annotations

import csv
import json
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


# Row-height model used by _compact_row_heights.
_TABLE_HEADER_HEIGHT_IN = 0.26
_TABLE_BODY_HEIGHT_IN = 0.21

# Maximum rows that fit per table container height.
EXEC_RECOMMENDATIONS_ROWS_ON_SLIDE = int((4.9 - _TABLE_HEADER_HEIGHT_IN) / _TABLE_BODY_HEIGHT_IN)
EXEC_ISSUES_ROWS_ON_SLIDE = int((0.95 - _TABLE_HEADER_HEIGHT_IN) / _TABLE_BODY_HEIGHT_IN)
TECH_CHECK_OUTCOMES_ROWS_ON_SLIDE = int((4.6 - _TABLE_HEADER_HEIGHT_IN) / _TABLE_BODY_HEIGHT_IN)
TECH_WATCHLIST_ROWS_ON_SLIDE = int((3.7 - _TABLE_HEADER_HEIGHT_IN) / _TABLE_BODY_HEIGHT_IN)
TECH_POLICY_EFFICACY_ROWS_ON_SLIDE = int((2.55 - _TABLE_HEADER_HEIGHT_IN) / _TABLE_BODY_HEIGHT_IN)
TECH_RECOMMENDATIONS_ROWS_ON_SLIDE = int((5.8 - _TABLE_HEADER_HEIGHT_IN) / _TABLE_BODY_HEIGHT_IN)

APPENDIX_ROWS_PER_SLIDE_STANDARD = int((5.8 - _TABLE_HEADER_HEIGHT_IN) / _TABLE_BODY_HEIGHT_IN)
APPENDIX_ROWS_PER_SLIDE_DENSE_TEXT = APPENDIX_ROWS_PER_SLIDE_STANDARD


def _safe_path_part(value: str) -> str:
    cleaned = "".join(ch for ch in value if ch.isalnum() or ch in ("-", "_", "."))
    return cleaned or "unknown_tenant"


def create_run_dir(output_root: Path, tenant_key: str) -> Path:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    tenant_dir = output_root / _safe_path_part(tenant_key)
    run_dir = tenant_dir / ts
    run_dir.mkdir(parents=True, exist_ok=True)
    return run_dir


def find_latest_run_dir(output_root: Path, tenant_key: str) -> Path | None:
    tenant_dir = output_root / _safe_path_part(tenant_key)
    if not tenant_dir.exists():
        return None
    run_dirs = [d for d in tenant_dir.iterdir() if d.is_dir()]
    if not run_dirs:
        return None
    run_dirs.sort(key=lambda p: p.name, reverse=True)
    return run_dirs[0]


def read_records_csv(csv_path: Path) -> list[dict[str, Any]]:
    if not csv_path.exists():
        return []
    with csv_path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle)
        return [dict(row) for row in reader]


def write_summary(run_dir: Path, summary: dict[str, Any]) -> Path:
    out_path = run_dir / "summary.json"
    out_path.write_text(json.dumps(summary, indent=2), encoding="utf-8")
    return out_path


def write_records_csv(run_dir: Path, filename: str, rows: list[dict[str, Any]]) -> Path | None:
    if not rows:
        return None
    out_path = run_dir / filename
    headers = sorted({k for row in rows for k in row.keys()})
    with out_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=headers)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)
    return out_path


def write_executive_markdown(run_dir: Path, summary: dict[str, Any]) -> Path:
    health = summary.get("health", {})
    checks = summary.get("checks", {})
    recs = summary.get("recommendations", [])

    lines: list[str] = []
    lines.append("# Carbon Black Health Check Executive Summary")
    lines.append("")
    lines.append(f"- Customer: {summary.get('customer_name', 'unknown')}")
    lines.append(f"- Tenant ID: {summary.get('tenant_id', 'unknown')}")
    lines.append(f"- Assessment Profile: {summary.get('assessment_profile', 'unknown')}")
    lines.append(f"- Score: {health.get('score', 'n/a')} ({health.get('status', 'n/a')})")
    lines.append("")

    lines.append("## Key Metrics")
    lines.append("")
    devices = checks.get("devices", {}).get("summary", {})
    alerts = checks.get("alerts", {}).get("summary", {})
    policy = checks.get("policy_posture", {}).get("summary", {})
    watchlists = checks.get("watchlists", {}).get("summary", {})
    drift = checks.get("policy_drift", {}).get("summary", {})
    policy_tuning = checks.get("policy_tuning_analysis", {}).get("summary", {})

    lines.append(f"- Devices: {devices.get('total_devices', 'n/a')} (active 7d ratio: {devices.get('active_ratio_last_7d', 'n/a')})")
    lines.append(f"- Alerts (30d): {alerts.get('total_alerts_30d', 'n/a')} (severity 7+: {alerts.get('high_severity_alerts', 'n/a')})")
    lines.append(f"- Policies: {policy.get('total_policies', 'n/a')} | Rules: {policy.get('total_rules', 'n/a')} | Blocking ratio: {policy.get('blocking_rule_ratio', 'n/a')}")
    lines.append(
        f"- Watchlists: {watchlists.get('total_watchlists', 'n/a')} "
        f"(alerting enabled: {watchlists.get('alerting_enabled_watchlists', 'n/a')}, "
        f"report count: {watchlists.get('total_watchlist_reports', 'n/a')})"
    )
    lines.append(
        f"- Policy drift: detected={drift.get('drift_detected', 'n/a')} "
        f"(changed={drift.get('changed_count', 'n/a')})"
    )
    drift_details = drift.get("changed_policy_details", [])
    if isinstance(drift_details, list) and drift_details:
        for detail in drift_details[:3]:
            policy_name = str(detail.get("policy_name", "unknown"))
            change_count = int(detail.get("change_count", 0))
            sample_changes = detail.get("changes", [])
            sample_fields = []
            if isinstance(sample_changes, list):
                sample_fields = [str(item.get("field", "")) for item in sample_changes[:3] if isinstance(item, dict)]
            lines.append(
                f"- Drift highlight: {policy_name} changed_fields={change_count} "
                f"sample={', '.join(sample_fields) if sample_fields else 'n/a'}"
            )
    lines.append("")

    lines.append("## Policy Analysis (Good/Better/Best)")
    lines.append("")
    if isinstance(policy_tuning, dict) and policy_tuning:
        framework = str(policy_tuning.get("framework", "Broadcom Endpoint Standard Good-Better-Best"))
        current_tier = str(policy_tuning.get("current_tier", "unknown"))
        score = policy_tuning.get("score_0_100", "n/a")
        applicable = bool(policy_tuning.get("framework_applicable", False))
        lines.append(f"- Framework: {framework}")
        lines.append(f"- Current tier: {current_tier}")
        lines.append(f"- Maturity score: {score}")
        lines.append(f"- Applicable: {'yes' if applicable else 'no'}")

        metrics = policy_tuning.get("metrics", {})
        if isinstance(metrics, dict):
            lines.append(
                "- Metrics: "
                f"fully_alert_only_policies={metrics.get('fully_alert_only_policies', 'n/a')}, "
                f"missing_required_controls={metrics.get('policies_with_missing_required_controls', 'n/a')}, "
                f"monitor_heavy_policy_groups={metrics.get('monitor_heavy_policy_groups', 'n/a')}, "
                f"permissions_p1={metrics.get('permissions_audit_p1_policies', 'n/a')}"
            )

        top_gaps = policy_tuning.get("top_gaps", [])
        if isinstance(top_gaps, list) and top_gaps:
            for gap in top_gaps[:4]:
                lines.append(f"- Gap: {gap}")

        next_actions = policy_tuning.get("next_actions", [])
        if isinstance(next_actions, list) and next_actions:
            for action in next_actions[:3]:
                lines.append(f"- Next action: {action}")
    else:
        lines.append("- Policy tuning analysis unavailable.")
    lines.append("")

    lines.append("## Prioritized Recommendations")
    lines.append("")
    for rec in recs:
        lines.append(
            f"- [{rec.get('priority', 'P4')}] {rec.get('area', 'General')}: "
            f"{rec.get('recommendation', '')} ({rec.get('evidence', '')})"
        )

    output_path = run_dir / "executive_summary.md"
    output_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return output_path


def _hex_color(value: str) -> RGBColor:
    return RGBColor(*bytes.fromhex(value))


def _get_check_summary(summary: dict[str, Any], check_name: str) -> dict[str, Any]:
    check = summary.get("checks", {}).get(check_name, {})
    if not isinstance(check, dict):
        return {}
    data = check.get("summary", {})
    return data if isinstance(data, dict) else {}


def _get_org_domain(summary: dict[str, Any]) -> str:
    checks = summary.get("checks", {}) if isinstance(summary, dict) else {}
    if isinstance(checks, dict):
        org_check = checks.get("org", {})
        if isinstance(org_check, dict):
            organization = org_check.get("organization", {})
            if isinstance(organization, dict):
                org_domain = organization.get("orgDomain")
                if org_domain is not None:
                    org_domain_str = str(org_domain).strip()
                    if org_domain_str:
                        return org_domain_str
    backend_url = str(summary.get("backend_url", "")).strip()
    return backend_url.replace("https://", "").replace("http://", "").rstrip("/")


def _org_defense_rules_disabled(summary: dict[str, Any]) -> bool:
    checks = summary.get("checks", {}) if isinstance(summary, dict) else {}
    if not isinstance(checks, dict):
        return False
    org_check = checks.get("org", {})
    if not isinstance(org_check, dict):
        return False
    organization = org_check.get("organization", {})
    if not isinstance(organization, dict):
        return False
    change_ts = organization.get("changeTimestamps", {})
    if not isinstance(change_ts, dict):
        return False
    timestamp_map = change_ts.get("timestampMap", {})
    if not isinstance(timestamp_map, dict):
        return False
    return "ORG_DISABLE_DEFENSE_RULES" in timestamp_map


def _format_ratio(value: Any) -> str:
    try:
        number = float(value)
    except Exception:
        return "n/a"
    if number <= 1:
        return f"{number:.1%}"
    return f"{number:.1f}%"


def _format_count(value: Any) -> str:
    try:
        return f"{int(value):,}"
    except Exception:
        return "n/a"


def _format_score(summary: dict[str, Any]) -> str:
    health = summary.get("health", {})
    if not isinstance(health, dict):
        return "n/a"
    score = health.get("score", "n/a")
    status = health.get("status", "n/a")
    return f"{score} / 100 ({status})"


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
    font_name: str = "Aptos",
    underline: bool = False,
    word_wrap: bool = True,
) -> Any:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.word_wrap = word_wrap
    frame.vertical_anchor = valign
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.text = text
    for para in frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = font_name
            run.font.underline = underline
            if color is not None:
                run.font.color.rgb = color
    return textbox


def _summarize_entitlements(entitlements: list[str] | Any, *, max_items: int | None = 2) -> str:
    if not isinstance(entitlements, list):
        return "auto-detect"
    products = [str(item).strip() for item in entitlements if str(item).strip()]
    if not products:
        return "auto-detect"
    if max_items is None or max_items <= 0:
        return ", ".join(products)
    shown = products[:max_items]
    remaining = len(products) - len(shown)
    if remaining > 0:
        return f"{', '.join(shown)} (+{remaining} more)"
    return ", ".join(shown)


def _add_background(slide, color: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_color(color)
    shape.line.fill.background()


def _add_unavailable_callout(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    subject: str,
    reason: str,
) -> None:
    _add_textbox(
        slide,
        left,
        top,
        width,
        height,
        f"{subject}\nUnavailable: {reason}.",
        font_size=11,
        color=_hex_color("5B6472"),
    )


def _add_eedr_not_applicable(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    subject: str,
) -> None:
    _add_unavailable_callout(slide, left, top, width, height, subject, "Not applicable for EEDR-only tenant")


def _add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    value: str,
    detail: str = "",
    *,
    accent: str = "1F3B73",
    value_font_size_max: int | None = None,
    detail_word_wrap: bool = False,
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _hex_color("FFFFFF")
    card.line.color.rgb = _hex_color("D7DEE8")
    card.line.width = Pt(1.1)

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.14),
        Inches(height),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _hex_color(accent)
    accent_bar.line.fill.background()

    _add_textbox(slide, left + 0.22, top + 0.15, width - 0.3, 0.28, title, font_size=11, bold=True, color=_hex_color("5B6472"))

    value_text = _safe_str(value)
    value_font_size = 24
    if height < 1.25:
        value_font_size = 20
    if height < 1.1:
        value_font_size = 18
    if len(value_text) > 14:
        value_font_size = min(value_font_size, 18)
    if len(value_text) > 18:
        value_font_size = min(value_font_size, 16)
    if len(value_text) > 22:
        value_font_size = min(value_font_size, 14)
    if len(value_text) > 26:
        value_font_size = min(value_font_size, 12)
    if isinstance(value_font_size_max, int) and value_font_size_max > 0:
        value_font_size = min(value_font_size, value_font_size_max)

    has_detail = bool(detail)
    if detail_word_wrap:
        detail_height = max(0.36, min(0.95, height * 0.38))
    elif has_detail:
        detail_height = 0.24 if height >= 1.2 else 0.2
    else:
        detail_height = 0.0

    detail_bottom_margin = 0.08
    detail_top = top + height - detail_height - detail_bottom_margin if has_detail else top + height

    value_top = top + 0.44
    value_bottom = detail_top - 0.06 if has_detail else top + height - 0.14
    value_height = max(0.22, value_bottom - value_top)
    _add_textbox(
        slide,
        left + 0.22,
        value_top,
        width - 0.3,
        value_height,
        value_text,
        font_size=value_font_size,
        bold=True,
        color=_hex_color("102033"),
        word_wrap=False,
    )
    if has_detail:
        detail_font_size = 10 if height >= 1.2 else 9
        _add_textbox(
            slide,
            left + 0.22,
            detail_top,
            width - 0.3,
            detail_height,
            detail,
            font_size=detail_font_size,
            color=_hex_color("5B6472"),
            word_wrap=detail_word_wrap,
        )


def _add_section_header(slide, title: str, subtitle: str = "", *, section_label: str = "SECTION") -> None:
    # Shared visual scaffold for all section slides across deck types.
    _ = section_label
    _add_textbox(slide, 0.55, 0.38, 9.6, 0.45, title, font_size=25, bold=True, color=_hex_color("102033"))
    if subtitle:
        _add_textbox(slide, 0.55, 0.8, 11.6, 0.3, subtitle, font_size=11, color=_hex_color("55606E"))

    top_rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(1.08),
        Inches(12.05),
        Inches(0.02),
    )
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = _hex_color("DCE5F1")
    top_rule.line.fill.background()


def _add_standard_cover_slide(
    prs: Presentation,
    *,
    deck_title: str,
    customer_name: str,
    tenant_key: str,
    assessment_profile: str,
    run_stamp: str,
    score_text: str,
    summary_text: str,
    right_cards: list[dict[str, str]],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "102033")
    _add_textbox(slide, 0.75, 0.62, 8.0, 1.0, deck_title, font_size=30, bold=True, color=_hex_color("FFFFFF"))
    _add_textbox(slide, 0.78, 1.64, 8.1, 0.38, customer_name, font_size=19, bold=True, color=_hex_color("DCE6F5"))
    _add_textbox(slide, 0.78, 2.02, 8.1, 0.3, f"Tenant {tenant_key}  |  Profile {assessment_profile}  |  Run {run_stamp}", font_size=11, color=_hex_color("B9C7D9"))
    _add_textbox(slide, 0.78, 2.42, 5.8, 0.3, "OVERALL SCORE", font_size=11, bold=True, color=_hex_color("93C5FD"))
    _add_textbox(slide, 0.78, 2.66, 8.6, 0.58, score_text, font_size=30, bold=True, color=_hex_color("7DD3FC"))
    _add_textbox(slide, 0.78, 3.38, 8.3, 0.62, summary_text, font_size=14, color=_hex_color("E5EEF8"))

    card_specs = [
        (right_cards[0], 0.82, 1.02),
        (right_cards[1], 1.98, 1.02),
        (right_cards[2], 3.14, 1.02),
        (right_cards[3], 4.3, 1.7),
    ]
    for card, top, height in card_specs:
        _add_card(
            slide,
            9.05,
            top,
            3.6,
            height,
            card.get("title", ""),
            card.get("value", ""),
            card.get("detail", ""),
            accent=card.get("accent", "60A5FA"),
            value_font_size_max=card.get("value_font_size_max"),
            detail_word_wrap=bool(card.get("detail_word_wrap", False)),
        )


def _add_appendix_divider_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "102033")
    _add_textbox(slide, 0.75, 0.86, 8.0, 0.56, "Appendix", font_size=36, bold=True, color=_hex_color("FFFFFF"))
    _add_textbox(
        slide,
        0.78,
        1.48,
        11.2,
        0.38,
        "Detailed reference tables and overflow data that support the core narrative.",
        font_size=13,
        color=_hex_color("DCE6F5"),
    )

    top_rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.75),
        Inches(2.0),
        Inches(11.95),
        Inches(0.03),
    )
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = _hex_color("334155")
    top_rule.line.fill.background()


def _add_bullet_list(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    font_size: int = 12,
    color: RGBColor | None = None,
    bullet_color: RGBColor | None = None,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.clear()

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {line}"
        paragraph.level = 0
        paragraph.bullet = False
        paragraph.space_after = Pt(7)
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = "Aptos"
            run.font.color.rgb = color if color is not None else _hex_color("23313F")
            if bullet_color is not None:
                paragraph.font.color.rgb = bullet_color


def _add_bold_label_bullet_list(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[tuple[str, str]],
    *,
    font_size: int = 12,
    color: RGBColor | None = None,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.clear()

    for index, (label, body) in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(7)

        label_run = paragraph.add_run()
        label_run.text = label
        label_run.font.size = Pt(font_size)
        label_run.font.name = "Aptos"
        label_run.font.bold = True
        label_run.font.color.rgb = color if color is not None else _hex_color("23313F")

        body_run = paragraph.add_run()
        body_run.text = f": {body}"
        body_run.font.size = Pt(font_size)
        body_run.font.name = "Aptos"
        body_run.font.color.rgb = color if color is not None else _hex_color("23313F")


def _top_items(breakdown: dict[str, Any], limit: int = 3) -> list[tuple[str, Any]]:
    if not isinstance(breakdown, dict):
        return []
    try:
        items = sorted(breakdown.items(), key=lambda item: int(item[1]), reverse=True)
    except Exception:
        items = list(breakdown.items())
    return items[:limit]


def _safe_str(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, (dict, list)):
        return json.dumps(value, ensure_ascii=True)
    return str(value)


def _truncate_chart_label(value: Any, max_length: int = 32) -> str:
    text = _safe_str(value).strip()
    if len(text) <= max_length:
        return text
    return text[: max_length - 3].rstrip() + "..."


def _appendix_rows_per_slide(headers: list[str]) -> int:
    _ = headers
    return APPENDIX_ROWS_PER_SLIDE_STANDARD


def _compact_row_heights(table: Any, max_height_in: float, data_row_count: int) -> float:
    header_height = 0.26
    body_height = 0.21

    if data_row_count <= 0:
        used_height = min(max_height_in, header_height)
        table.rows[0].height = Inches(used_height)
        return used_height

    target_height = header_height + (body_height * data_row_count)
    used_height = min(max_height_in, target_height)

    min_body_height = 0.14
    max_header_height = max(used_height - (min_body_height * data_row_count), 0.16)
    actual_header_height = min(header_height, max_header_height)
    remaining = max(used_height - actual_header_height, min_body_height * data_row_count)
    actual_body_height = remaining / data_row_count

    table.rows[0].height = Inches(actual_header_height)
    for row_idx in range(1, data_row_count + 1):
        table.rows[row_idx].height = Inches(actual_body_height)

    return used_height


def _fit_column_widths(table: Any, headers: list[str], rows: list[list[Any]], table_width_in: float) -> None:
    if not headers:
        return

    # Estimate relative width from visible text density, then enforce a minimum width.
    weights: list[float] = []
    for col_idx, header in enumerate(headers):
        max_len = len(_safe_str(header).strip())
        for row in rows:
            if col_idx >= len(row):
                continue
            value_len = len(_safe_str(row[col_idx]).strip())
            if value_len > max_len:
                max_len = value_len
        # Cap influence of extreme values and keep narrow columns from collapsing.
        weights.append(max(6.0, min(float(max_len), 42.0)))

    weight_total = sum(weights) or float(len(headers))
    min_col_width_in = 0.8
    widths = [(table_width_in * (weight / weight_total)) for weight in weights]

    deficit = 0.0
    for idx, width_in in enumerate(widths):
        if width_in < min_col_width_in:
            deficit += (min_col_width_in - width_in)
            widths[idx] = min_col_width_in

    if deficit > 0:
        flex_indices = [idx for idx, width_in in enumerate(widths) if width_in > min_col_width_in]
        while deficit > 0 and flex_indices:
            flex_total = sum(widths[idx] - min_col_width_in for idx in flex_indices)
            if flex_total <= 0:
                break
            for idx in list(flex_indices):
                flex = widths[idx] - min_col_width_in
                share = deficit * (flex / flex_total)
                reduction = min(share, flex)
                widths[idx] -= reduction
                deficit -= reduction
                if widths[idx] <= min_col_width_in + 1e-6:
                    widths[idx] = min_col_width_in
            flex_indices = [idx for idx, width_in in enumerate(widths) if width_in > min_col_width_in]

    for col_idx, width_in in enumerate(widths):
        table.columns[col_idx].width = Inches(width_in)


def _add_table(slide, left: float, top: float, width: float, height: float, headers: list[str], rows: list[list[Any]]) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    _fit_column_widths(table, headers, rows, width)
    used_height = _compact_row_heights(table, height, len(rows))
    table_shape.height = Inches(used_height)

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.margin_left = Pt(3)
        cell.margin_right = Pt(3)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            if paragraph.runs:
                run = paragraph.runs[0]
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.name = "Aptos"
                run.font.color.rgb = _hex_color("FFFFFF")
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_color("102033")

    for row_index, row_values in enumerate(rows, start=1):
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.text = _safe_str(value)
            cell.margin_left = Pt(3)
            cell.margin_right = Pt(3)
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                if col_index == 0:
                    paragraph.alignment = PP_ALIGN.LEFT
                else:
                    paragraph.alignment = PP_ALIGN.CENTER
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)
                if paragraph.runs:
                    run = paragraph.runs[0]
                    run.font.size = Pt(9)
                    run.font.name = "Aptos"
                    run.font.color.rgb = _hex_color("23313F")
            if row_index % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_color("F4F7FB")


def _to_numeric(value: Any) -> float:
    try:
        return float(value)
    except Exception:
        return 0.0


def _sorted_mapping_items(mapping: dict[str, Any], limit: int = 8, *, sort_by_value: bool = True) -> list[tuple[str, float]]:
    if not isinstance(mapping, dict):
        return []
    items: list[tuple[str, float]] = [(str(key), _to_numeric(val)) for key, val in mapping.items()]
    if sort_by_value:
        items.sort(key=lambda item: item[1], reverse=True)
    else:
        items.sort(key=lambda item: item[0])
    return items[:limit]


def _apply_chart_theme(
    chart: Any,
    title: str,
    *,
    show_legend: bool = False,
    legend_position: XL_LEGEND_POSITION = XL_LEGEND_POSITION.BOTTOM,
    category_font_size: int = 9,
    value_font_size: int = 9,
) -> None:
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    if chart.chart_title.text_frame.paragraphs:
        p = chart.chart_title.text_frame.paragraphs[0]
        if p.runs:
            run = p.runs[0]
            run.font.name = "Aptos"
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = _hex_color("102033")

    chart.has_legend = show_legend
    if show_legend and chart.legend is not None:
        chart.legend.position = legend_position
        chart.legend.include_in_layout = False
        chart.legend.font.name = "Aptos"
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = _hex_color("374151")

    try:
        category_axis = chart.category_axis
        category_axis.tick_labels.font.name = "Aptos"
        category_axis.tick_labels.font.size = Pt(category_font_size)
        category_axis.tick_labels.font.color.rgb = _hex_color("374151")
    except Exception:
        pass

    try:
        value_axis = chart.value_axis
        value_axis.tick_labels.font.name = "Aptos"
        value_axis.tick_labels.font.size = Pt(value_font_size)
        value_axis.tick_labels.font.color.rgb = _hex_color("374151")
    except Exception:
        pass


def _style_chart_data_labels(labels: Any, font_size: int = 8) -> None:
    labels.font.name = "Aptos"
    labels.font.size = Pt(font_size)
    labels.font.color.rgb = _hex_color("1F2937")


def _force_all_category_tick_labels(chart: Any) -> None:
    # PowerPoint can auto-skip category labels on dense bar charts.
    # Set the XML skip markers directly so every category is eligible for display.
    try:
        category_axis = chart.category_axis
        category_axis.tick_label_spacing = 1
    except Exception:
        pass

    try:
        cat_ax_elements = chart._chartSpace.xpath(".//c:catAx")
    except Exception:
        cat_ax_elements = []

    for cat_ax in cat_ax_elements:
        for tag_name in ("c:tickLblSkip", "c:tickMarkSkip"):
            try:
                existing_nodes = cat_ax.xpath(f"./{tag_name}")
                if existing_nodes:
                    existing = existing_nodes[0]
                else:
                    existing = OxmlElement(tag_name)
                    cat_ax.append(existing)
                existing.set("val", "1")
            except Exception:
                continue


def _add_column_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    categories: list[str],
    values: list[float],
    *,
    series_name: str = "Count",
    show_data_labels: bool = True,
) -> None:
    if not categories or not values:
        _add_unavailable_callout(slide, left, top, width, height, title, "No chart data available")
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    _apply_chart_theme(chart, title, show_legend=False, category_font_size=9, value_font_size=9)
    chart.value_axis.has_major_gridlines = True
    if show_data_labels:
        chart.plots[0].has_data_labels = True
        _style_chart_data_labels(chart.plots[0].data_labels, font_size=8)


def _add_bar_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    categories: list[str],
    values: list[float],
    *,
    series_name: str = "Count",
    show_data_labels: bool = True,
    category_font_size: int = 10,
    value_font_size: int = 9,
) -> None:
    if not categories or not values:
        _add_unavailable_callout(slide, left, top, width, height, title, "No chart data available")
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    _apply_chart_theme(
        chart,
        title,
        show_legend=False,
        category_font_size=category_font_size,
        value_font_size=value_font_size,
    )
    _force_all_category_tick_labels(chart)
    chart.value_axis.has_major_gridlines = True
    if show_data_labels:
        chart.plots[0].has_data_labels = True
        _style_chart_data_labels(chart.plots[0].data_labels, font_size=9)


def _add_doughnut_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    categories: list[str],
    values: list[float],
    *,
    series_name: str = "Share",
) -> None:
    if not categories or not values:
        _add_unavailable_callout(slide, left, top, width, height, title, "No chart data available")
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    _apply_chart_theme(chart, title, show_legend=True, legend_position=XL_LEGEND_POSITION.RIGHT)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.show_percentage = True
    _style_chart_data_labels(chart.plots[0].data_labels, font_size=8)


def _add_pie_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    categories: list[str],
    values: list[float],
    *,
    series_name: str = "Share",
) -> None:
    if not categories or not values or sum(values) <= 0:
        _add_unavailable_callout(slide, left, top, width, height, title, "No chart data available")
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    _apply_chart_theme(chart, title, show_legend=True, legend_position=XL_LEGEND_POSITION.RIGHT)
    # Keep the pie clean: rely on legend-only categories to avoid label overlap.
    chart.plots[0].has_data_labels = False


def _add_line_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    categories: list[str],
    values: list[float],
    *,
    series_name: str = "Trend",
) -> None:
    if not categories or not values:
        _add_unavailable_callout(slide, left, top, width, height, title, "No chart data available")
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart
    _apply_chart_theme(chart, title, show_legend=False, category_font_size=8, value_font_size=9)
    chart.value_axis.has_major_gridlines = True


def _daily_alert_trend_data(daily_counts: dict[str, Any], max_points: int = 45) -> tuple[list[str], list[float]]:
    items = _sorted_mapping_items(daily_counts, limit=500, sort_by_value=False)
    if max_points > 0 and len(items) > max_points:
        items = items[-max_points:]
    categories = [key[5:] if len(key) >= 10 else key for key, _ in items]
    values = [value for _, value in items]
    return categories, values


def _add_live_query_audit_slide(prs: Presentation, live_query: dict[str, Any]) -> None:
    """Render Live Query (Audit and Remediation) analytics slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(
        slide,
        "Live Query (Audit and Remediation)",
        "Who ran queries, what was queried, and how activity trends over time.",
    )

    total_events = int(_to_numeric(live_query.get("total_live_query_events", 0)))
    total_runs = int(_to_numeric(live_query.get("total_query_runs", 0)))
    recommended_runs = int(_to_numeric(live_query.get("recommended_query_runs", 0)))
    custom_runs = int(_to_numeric(live_query.get("custom_query_runs", 0)))
    avg_endpoints_raw = live_query.get("avg_endpoints_per_query")
    if isinstance(avg_endpoints_raw, (int, float)):
        avg_endpoints = f"{avg_endpoints_raw:.2f}"
    else:
        avg_endpoints = "n/a"

    _add_card(slide, 0.65, 1.2, 2.45, 1.2, "Live Query Events", _format_count(total_events), "audit log entries", accent="2563EB")
    _add_card(slide, 3.25, 1.2, 2.45, 1.2, "Query Runs", _format_count(total_runs), "run actions", accent="0F766E")
    _add_card(slide, 5.85, 1.2, 2.55, 1.2, "Recommended", _format_count(recommended_runs), "runs", accent="22C55E")
    _add_card(slide, 8.55, 1.2, 2.55, 1.2, "Custom", _format_count(custom_runs), "runs", accent="DC2626")

    creators = live_query.get("users_query_creates", []) if isinstance(live_query, dict) else []
    creator_rows: list[list[Any]] = []
    if isinstance(creators, list):
        for entry in creators[:10]:
            if not isinstance(entry, dict):
                continue
            creator_rows.append([
                str(entry.get("user", "unknown")),
                int(_to_numeric(entry.get("run_count", entry.get("create_count", 0)))),
            ])
    _add_table(
        slide,
        0.65,
        2.65,
        3.95,
        2.45,
        ["User", "Runs"],
        creator_rows or [["No run events", 0]],
    )

    top_queries = live_query.get("top_queries", []) if isinstance(live_query, dict) else []
    query_rows: list[list[Any]] = []
    if isinstance(top_queries, list):
        for entry in top_queries[:8]:
            if not isinstance(entry, dict):
                continue
            query_rows.append([
                str(entry.get("query_name", "Unknown")),
                int(_to_numeric(entry.get("run_count", 0))),
                str(entry.get("query_type", "Unknown")),
                str(entry.get("os", "Unknown")),
            ])
    _add_table(
        slide,
        4.8,
        2.65,
        7.85,
        2.45,
        ["Query", "Runs", "Type", "OS"],
        query_rows or [["No query run events", 0, "n/a", "n/a"]],
    )

    os_breakdown = live_query.get("queried_os_breakdown", {}) if isinstance(live_query, dict) else {}
    os_items = _sorted_mapping_items(os_breakdown, limit=5, sort_by_value=True)
    _add_doughnut_chart(
        slide,
        0.65,
        5.2,
        4.6,
        2.0,
        f"Queried OS Mix (avg endpoints/query: {avg_endpoints})",
        [name for name, _ in os_items],
        [value for _, value in os_items],
        series_name="Runs",
    )

    daily_counts = live_query.get("daily_event_counts", {}) if isinstance(live_query, dict) else {}
    trend_labels, trend_values = _daily_alert_trend_data(daily_counts, max_points=30)
    _add_line_chart(
        slide,
        5.4,
        5.2,
        7.25,
        2.0,
        "Live Query Event Timeline",
        trend_labels,
        trend_values,
        series_name="Events/day",
    )

    _add_textbox(
        slide,
        0.65,
        7.16,
        12.0,
        0.2,
        "Recommended vs custom classification uses recommended_query_id (non-null => Recommended).",
        font_size=8,
        color=_hex_color("55606E"),
    )


def _rows_from_mapping(mapping: dict[str, Any], limit: int = 8) -> list[list[Any]]:
    if not isinstance(mapping, dict):
        return []
    try:
        items = sorted(mapping.items(), key=lambda item: item[0])
    except Exception:
        items = list(mapping.items())
    return [[key, value] for key, value in items[:limit]]


def _policy_settings_concern_items(policy_settings: dict[str, Any]) -> list[tuple[str, float]]:
    ordered_concern_keys = [
        "submit_unknown_binaries_enabled",
        "live_response_disabled",
        "delay_execute_for_cloud_scan_disabled",
        "scan_execution_on_network_drive_disabled",
        "scan_network_drive_enabled",
        "uninstall_code_disabled",
        "auth_event_collection_disabled",
        "xdr_network_data_collection_disabled",
        "windows_security_center_disabled",
    ]

    concerns = policy_settings.get("settings_of_concern", {}) if isinstance(policy_settings, dict) else {}
    if isinstance(concerns, dict) and concerns:
        return [(key, _to_numeric(concerns.get(key, 0))) for key in ordered_concern_keys]

    # Backward-compatible fallback for historical summaries without settings_of_concern.
    settings = policy_settings.get("settings", {}) if isinstance(policy_settings, dict) else {}
    if not isinstance(settings, dict):
        return []

    fallback_rules = {
        "submit_unknown_binaries_enabled": ("submit_unknown_binaries", "true"),
        "live_response_disabled": ("live_response_enabled", "false"),
        "delay_execute_for_cloud_scan_disabled": ("delay_execute_for_cloud_scan", "false"),
        "scan_execution_on_network_drive_disabled": ("scan_execution_on_network_drive", "false"),
        "scan_network_drive_enabled": ("scan_network_drive", "true"),
        "uninstall_code_disabled": ("uninstall_code_enabled", "false"),
        "auth_event_collection_disabled": ("auth_event_collection", "false"),
        "xdr_network_data_collection_disabled": ("xdr_network_data_collection", "false"),
        "windows_security_center_disabled": ("windows_security_center", "false"),
    }
    rows: list[tuple[str, float]] = []
    for concern_key, (setting_key, concern_state) in fallback_rules.items():
        counts = settings.get(setting_key, {}) if isinstance(settings.get(setting_key, {}), dict) else {}
        rows.append((concern_key, _to_numeric(counts.get(concern_state, 0))))
    return rows


def _policy_setting_display_name(concern_key: str) -> str:
    labels = {
        "submit_unknown_binaries_enabled": "Submit Unknown Binaries Enabled",
        "live_response_disabled": "Live Response Disabled",
        "delay_execute_for_cloud_scan_disabled": "Delay Execute for Cloud Scan Disabled",
        "scan_execution_on_network_drive_disabled": "Scan Execution on Network Drive Disabled",
        "scan_network_drive_enabled": "Scan Network Drive Enabled",
        "uninstall_code_disabled": "Uninstall Code Disabled",
        "auth_event_collection_disabled": "Auth Event Collection Disabled",
        "xdr_network_data_collection_disabled": "XDR Network Data Collection Disabled",
        "windows_security_center_disabled": "Windows Security Center Disabled",
    }
    return labels.get(concern_key, concern_key.replace("_", " ").title())


def _policy_setting_rationale(concern_key: str) -> str:
    rationales = {
        "submit_unknown_binaries_enabled": "Uploads a first-seen binary to secure tenant storage so analysts can retrieve it later, even if the file is deleted. Organizations with strict proprietary-code controls may prefer exceptions or disabling this setting.",
        "live_response_disabled": "Removes remote response capability, limiting containment and host-level investigation after a detection.",
        "delay_execute_for_cloud_scan_disabled": "Allows files to run before the cloud verdict is returned, reducing the execution-time safety check.",
        "scan_execution_on_network_drive_disabled": "Leaves executables on network drives unchecked at launch, creating a security gap. This control is usually lower impact than full network-drive scanning because it evaluates files only when they run.",
        "scan_network_drive_enabled": "Can add noticeable performance overhead on large or busy network shares because more remote content is scanned.",
        "uninstall_code_disabled": "Makes sensor removal easier because endpoints do not require an uninstall code.",
        "auth_event_collection_disabled": "Reduces identity telemetry available for investigation and cross-signal correlation.",
        "xdr_network_data_collection_disabled": "Reduces network telemetry available to XDR, weakening visibility and investigation context.",
        "windows_security_center_disabled": "Reduces Windows Security Center integration, making endpoint health and AV status less visible at the OS layer.",
    }
    return rationales.get(concern_key, "Review this setting to confirm it matches the intended security baseline.")


def _rows_from_list(items: Any, columns: list[str], limit: int | None = None) -> list[list[Any]]:
    if not isinstance(items, list):
        return []
    rows: list[list[Any]] = []
    selected_items = items if limit is None else items[:limit]
    for item in selected_items:
        if not isinstance(item, dict):
            continue
        rows.append([item.get(column, "") for column in columns])
    return rows


def _watchlist_detail_rows(items: Any) -> list[list[Any]]:
    if not isinstance(items, list):
        return []

    normalized: list[tuple[int, int, str, list[Any]]] = []
    for item in items:
        if not isinstance(item, dict):
            continue

        name = _safe_str(item.get("name", "unknown")) or "unknown"
        enabled_raw = item.get("enabled")
        alerting_raw = item.get("alerting_enabled")
        report_count = item.get("report_count")

        enabled_text = "Yes" if enabled_raw is True else ("No" if enabled_raw is False else "n/a")
        alerting_text = "Yes" if alerting_raw is True else ("No" if alerting_raw is False else "n/a")
        report_text = "n/a" if report_count in (None, "") else _safe_str(report_count)

        if isinstance(report_count, (int, float)):
            report_sort = int(report_count)
        else:
            report_sort = -1

        alert_rank = 0 if alerting_text == "Yes" else 1 if alerting_text == "No" else 2
        normalized.append((alert_rank, -report_sort, name.lower(), [name, enabled_text, alerting_text, report_text]))

    normalized.sort(key=lambda item: (item[0], item[1], item[2]))
    return [row for _, _, _, row in normalized]


def _add_sampling_footnote(slide, total: int, displayed: int, y_position: float = 7.2, always_show: bool = False) -> Any | None:
    """Add a footnote indicating data context (completeness, pagination, etc)."""
    if not always_show and total <= displayed:
        return None
    
    if total <= displayed:
        text = f"Showing all {total} items."
        footnote_color = _hex_color("6B7280")
        underline = False
    else:
        ratio = (displayed / total) * 100
        text = f"Showing {displayed} of {total} items ({ratio:.0f}%). Click to open appendix full list."
        footnote_color = _hex_color("2563EB")
        underline = True
    
    return _add_textbox(
        slide,
        0.55,
        y_position,
        12.2,
        0.22,
        text,
        font_size=8,
        color=footnote_color,
        underline=underline,
    )


def _dormant_recommendation_appendix_data(
    recommendations: Any,
    user_logins: dict[str, Any],
    *,
    displayed_names_in_evidence: int = 5,
) -> tuple[list[list[Any]], list[str], int, int, str] | None:
    if not isinstance(recommendations, list) or not isinstance(user_logins, dict):
        return None

    for rec in recommendations:
        if not isinstance(rec, dict):
            continue
        if str(rec.get("area", "")).strip().lower() != "user logins":
            continue
        evidence = str(rec.get("evidence", ""))
        if "dormant_over_60d_count=" in evidence:
            dormant_details = user_logins.get("dormant_over_60d_details", [])
            if isinstance(dormant_details, list) and len(dormant_details) > displayed_names_in_evidence:
                rows = [
                    [
                        "60d+",
                        str(item.get("user", "unknown")) if isinstance(item, dict) else str(item),
                        str(item.get("role", "unknown")) if isinstance(item, dict) else "unknown",
                    ]
                    for item in dormant_details
                ]
                rows.sort(key=lambda row: (str(row[2]).lower(), str(row[1]).lower()))
                return rows, ["Dormancy Bucket", "User", "Role"], displayed_names_in_evidence, len(dormant_details), "Recommendation Dormant Users"

            dormant_users = user_logins.get("dormant_over_60d", [])
            if isinstance(dormant_users, list) and len(dormant_users) > displayed_names_in_evidence:
                rows = [["60d+", user, "unknown"] for user in dormant_users]
                return rows, ["Dormancy Bucket", "User", "Role"], displayed_names_in_evidence, len(dormant_users), "Recommendation Dormant Users"
            return None
        if "dormant_over_30d_count=" in evidence:
            dormant_details = user_logins.get("dormant_over_30d_details", [])
            if isinstance(dormant_details, list) and len(dormant_details) > displayed_names_in_evidence:
                rows = [
                    [
                        "30d+",
                        str(item.get("user", "unknown")) if isinstance(item, dict) else str(item),
                        str(item.get("role", "unknown")) if isinstance(item, dict) else "unknown",
                    ]
                    for item in dormant_details
                ]
                rows.sort(key=lambda row: (str(row[2]).lower(), str(row[1]).lower()))
                return rows, ["Dormancy Bucket", "User", "Role"], displayed_names_in_evidence, len(dormant_details), "Recommendation Dormant Users"

            dormant_users = user_logins.get("dormant_over_30d", [])
            if isinstance(dormant_users, list) and len(dormant_users) > displayed_names_in_evidence:
                rows = [["30d+", user, "unknown"] for user in dormant_users]
                return rows, ["Dormancy Bucket", "User", "Role"], displayed_names_in_evidence, len(dormant_users), "Recommendation Dormant Users"
            return None

    return None


def _is_policy_recommendation(rec: dict[str, Any]) -> bool:
    area = str(rec.get("area", "")).strip().lower()
    if not area:
        return False
    if "policy" in area:
        return True
    return area in {"core prevention", "permissions rule audit"}


def _recommendation_rows(recommendations: Any, *, policy_only: bool | None = None) -> list[list[Any]]:
    if not isinstance(recommendations, list):
        return []

    rows: list[list[Any]] = []
    for rec in recommendations:
        if not isinstance(rec, dict):
            continue
        is_policy = _is_policy_recommendation(rec)
        if policy_only is True and not is_policy:
            continue
        if policy_only is False and is_policy:
            continue
        rows.append(
            [
                rec.get("priority", "P4"),
                rec.get("area", "General"),
                rec.get("recommendation", ""),
                rec.get("evidence", ""),
            ]
        )
    return rows


def _policy_summary_row(policy_tuning: dict[str, Any], policy_rec_count: int) -> list[Any] | None:
    if not isinstance(policy_tuning, dict) or policy_rec_count <= 0:
        return None

    current_tier = str(policy_tuning.get("current_tier", "unknown"))
    score = policy_tuning.get("score_0_100", "n/a")
    top_gaps = policy_tuning.get("top_gaps", [])
    if isinstance(top_gaps, list) and top_gaps:
        gap_hint = str(top_gaps[0])
    else:
        gap_hint = "See appendix for detailed policy actions"

    return [
        "P2",
        "Policy Summary",
        f"Policy maturity tier is {current_tier} (score {score}/100). Detailed policy recommendations are in the appendix.",
        f"policy_recommendations={policy_rec_count}; top_gap={gap_hint}",
    ]


def _policy_maturity_action_plan_items(policy_tuning: dict[str, Any], *, attention_only: bool = True) -> list[dict[str, Any]]:
    if not isinstance(policy_tuning, dict):
        return []

    key = "policy_maturity_action_plan" if attention_only else "policy_maturity_action_plan_all"
    plan = policy_tuning.get(key, [])
    if not isinstance(plan, list):
        return []

    items: list[dict[str, Any]] = []
    for item in plan:
        if not isinstance(item, dict):
            continue
        actions_raw = item.get("actions", [])
        if isinstance(actions_raw, list):
            actions = [str(action).strip() for action in actions_raw if str(action).strip()]
        elif str(actions_raw).strip():
            actions = [str(actions_raw).strip()]
        else:
            actions = []
        items.append(
            {
                "policy_name": str(item.get("policy_name", "unknown")),
                "assigned_endpoints": int(_to_numeric(item.get("assigned_endpoints", 0))),
                "current_tier": str(item.get("current_tier", "unknown")).upper(),
                "next_target": str(item.get("next_target", "maintain")).upper(),
                "actions": actions,
            }
        )
    return items


def _policy_maturity_action_rows(policy_tuning: dict[str, Any], *, attention_only: bool = True) -> list[list[Any]]:
    items = _policy_maturity_action_plan_items(policy_tuning, attention_only=attention_only)

    rows: list[list[Any]] = []
    for item in items:
        actions_text = "; ".join(item.get("actions", []))
        rows.append(
            [
                str(item.get("policy_name", "unknown")),
                int(_to_numeric(item.get("assigned_endpoints", 0))),
                str(item.get("current_tier", "unknown")).upper(),
                str(item.get("next_target", "maintain")).upper(),
                actions_text or "No action required.",
            ]
        )
    return rows


def _add_policy_maturity_action_block(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    item: dict[str, Any],
) -> None:
    block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    block.fill.solid()
    block.fill.fore_color.rgb = _hex_color("F8FAFC")
    block.line.color.rgb = _hex_color("CBD5E1")
    block.line.width = Pt(1)

    policy_name = str(item.get("policy_name", "unknown"))
    assigned_endpoints = _format_count(item.get("assigned_endpoints", 0))
    current_tier = str(item.get("current_tier", "UNKNOWN"))
    next_target = str(item.get("next_target", "MAINTAIN"))
    actions = item.get("actions", []) if isinstance(item.get("actions", []), list) else []
    actions_to_show = [str(action).strip() for action in actions if str(action).strip()][:4]
    if not actions_to_show:
        actions_to_show = ["No action required."]

    _add_textbox(
        slide,
        left + 0.18,
        top + 0.12,
        width - 0.36,
        0.32,
        policy_name,
        font_size=14,
        bold=True,
        color=_hex_color("102033"),
    )
    _add_textbox(
        slide,
        left + 0.18,
        top + 0.48,
        width - 0.36,
        0.26,
        f"Assigned endpoints: {assigned_endpoints}   |   {current_tier} -> {next_target}",
        font_size=10,
        color=_hex_color("475569"),
    )
    _add_bullet_list(
        slide,
        left + 0.18,
        top + 0.82,
        width - 0.36,
        height - 0.94,
        actions_to_show,
        font_size=10,
        color=_hex_color("23313F"),
    )


def _paginate_policy_maturity_action_plan(
    prs: Presentation,
    items: list[dict[str, Any]],
    title: str,
    items_per_slide: int = 2,
) -> int:
    if not items:
        return 0

    slides_created = 0

    summary_rows: list[list[Any]] = []
    summary_items = sorted(
        items,
        key=lambda item: (
            -int(_to_numeric(item.get("assigned_endpoints", 0))),
            str(item.get("policy_name", "")).lower(),
        ),
    )
    for item in summary_items:
        summary_rows.append(
            [
                str(item.get("policy_name", "unknown")),
                str(item.get("current_tier", "UNKNOWN")).upper(),
                int(_to_numeric(item.get("assigned_endpoints", 0))),
            ]
        )
    slides_created += _paginate_table_to_slides(
        prs,
        summary_rows,
        ["Policy", "Maturity Tier", "Assigned Endpoints"],
        f"{title} - Summary",
        rows_per_slide=max(20, len(summary_rows)),
    )

    for page_idx in range(0, len(items), items_per_slide):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_background(slide, "FFFFFF")
        page_num = (page_idx // items_per_slide) + 1
        total_pages = (len(items) + items_per_slide - 1) // items_per_slide
        page_title = f"{title} (Page {page_num}/{total_pages})"
        _add_section_header(
            slide,
            page_title,
            "Action-plan appendix view: per-policy next step bullets are shown instead of a compressed table.",
        )

        page_items = items[page_idx : page_idx + items_per_slide]
        positions = [
            (0.65, 1.45),
            (0.65, 4.1),
        ]
        for item, (left, top) in zip(page_items, positions):
            _add_policy_maturity_action_block(slide, left, top, 12.0, 2.5, item)
        slides_created += 1

    return slides_created


def _paginate_table_to_slides(
    prs: Presentation,
    all_rows: list[list[Any]],
    headers: list[str],
    title: str,
    rows_per_slide: int = APPENDIX_ROWS_PER_SLIDE_STANDARD,
) -> int:
    """Create paginated slides for a large table. Returns number of slides created."""
    if not all_rows:
        return 0

    slides_created = 0
    for page_idx in range(0, len(all_rows), rows_per_slide):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_background(slide, "FFFFFF")
        page_num = (page_idx // rows_per_slide) + 1
        total_pages = (len(all_rows) + rows_per_slide - 1) // rows_per_slide
        page_title = f"{title} (Page {page_num}/{total_pages})"
        _add_section_header(slide, page_title, "")

        page_rows = all_rows[page_idx : page_idx + rows_per_slide]
        _add_table(slide, 0.55, 1.3, 12.2, 5.8, headers, page_rows)
        slides_created += 1

    return slides_created


def _find_slide_by_title_prefix(prs: Presentation, title_prefix: str) -> Any | None:
    prefix = (title_prefix or "").strip().lower()
    if not prefix:
        return None

    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip().lower()
            if text.startswith(prefix):
                return slide
    return None


def _add_score_breakdown_slide(prs: Presentation, summary: dict[str, Any], assessment_profile: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "Score Breakdown", "Overall score is an internal weighted model based on key risk signals in this run.")

    health_data = summary.get("health", {}) if isinstance(summary, dict) else {}
    current_score = _to_numeric(health_data.get("score", 0))
    current_status = str(health_data.get("status", "n/a"))
    current_profile = str(health_data.get("assessment_profile", assessment_profile))

    _add_card(
        slide,
        0.65,
        1.25,
        3.2,
        1.35,
        "Current Score",
        _format_count(current_score),
        f"Status {current_status}",
        accent="111827",
    )
    _add_card(
        slide,
        3.95,
        1.25,
        3.2,
        1.35,
        "Scoring Model",
        "Weighted",
        f"Profile {current_profile}",
        accent="2563EB",
    )
    _add_card(
        slide,
        7.25,
        1.25,
        2.65,
        1.35,
        "Good",
        "85-100",
        "low residual risk",
        accent="22C55E",
    )
    _add_card(
        slide,
        10.0,
        1.25,
        2.65,
        1.35,
        "Watch",
        "65-84",
        "needs attention",
        accent="F59E0B",
    )

    score_range_rows = [
        ["85-100", "good", "healthy posture with manageable findings"],
        ["65-84", "watch", "material gaps exist; prioritize remediation"],
        ["0-64", "at_risk", "high-risk posture requiring immediate action"],
    ]
    _add_table(
        slide,
        0.65,
        2.9,
        6.9,
        1.95,
        ["Score Range", "Status", "Interpretation"],
        score_range_rows,
    )

    _add_column_chart(
        slide,
        7.8,
        2.9,
        4.9,
        1.95,
        "Status Thresholds",
        ["At Risk", "Watch", "Good"],
        [64, 84, 100],
        series_name="Upper bound",
    )

    score_notes = [
        "Overall score starts at 100 and applies weighted deductions.",
        "Primary drivers: sensor activity, alert severity mix (7+ and sub-5), policy posture, watchlist effectiveness, and policy maturity/tuning.",
    ]
    health_notes = health_data.get("notes", []) if isinstance(health_data, dict) else []
    if isinstance(health_notes, list):
        for note in health_notes[:2]:
            score_notes.append(f"Current run note: {note}")

    _add_bullet_list(slide, 0.8, 5.2, 11.9, 1.8, score_notes[:4])


def _add_metric_definitions_slide(
    prs: Presentation,
    *,
    sensor_quality: dict[str, Any],
    alert_quality: dict[str, Any],
    policy_efficacy: dict[str, Any],
) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(
        slide,
        "Metric Definitions",
        "These definitions explain how to read the quality and efficacy signals shown throughout the deck.",
    )

    _add_card(
        slide,
        0.65,
        1.35,
        2.95,
        1.15,
        "Sensor Quality",
        _format_count(sensor_quality.get("stale_sensors_over_7d")),
        "stale sensors >7d",
        accent="F59E0B",
    )
    _add_card(
        slide,
        3.75,
        1.35,
        2.95,
        1.15,
        "Alert Quality",
        _format_count(alert_quality.get("repeated_detection_keys_over_3")),
        "repeated alert signatures >3",
        accent="DC2626",
    )
    _add_card(
        slide,
        6.85,
        1.35,
        2.95,
        1.15,
        "Efficacy",
        _format_count(policy_efficacy.get("total_policy_groups")),
        "policy groups scored",
        accent="111827",
    )
    _add_card(
        slide,
        9.95,
        1.35,
        2.7,
        1.15,
        "Alert Noise",
        _format_ratio(alert_quality.get("noise_ratio")),
        "lower is better",
        accent="7C3AED",
    )

    definition_rows = [
        [
            "Sensor Quality",
            "Endpoint telemetry hygiene using stale sensors and version health.",
            "High stale or unknown versions mean reduced endpoint visibility.",
        ],
        [
            "Alert Quality",
            "Actionability of alert stream using repeated alert signatures and aging buckets.",
            "More duplicates and older unresolved alerts indicate triage friction.",
        ],
        [
            "Efficacy",
            "Policy enforcement strength based on monitor-vs-block rule mix by policy group.",
            "Higher effective block posture generally means stronger prevention outcomes.",
        ],
        [
            "Alert Noise",
            "Low-signal share of alerts represented by noise_ratio.",
            "Lower noise ratio means analysts can focus on higher-confidence events.",
        ],
    ]
    _add_table(
        slide,
        0.65,
        2.75,
        12.0,
        3.75,
        ["Metric", "What It Means", "How To Read It"],
        definition_rows,
    )

    return slide


def write_executive_pptx(run_dir: Path) -> Path:
    summary_path = run_dir / "summary.json"
    if not summary_path.exists():
        raise FileNotFoundError(f"summary.json not found in {run_dir}")

    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    org_domain = _safe_path_part(_get_org_domain(summary))
    pptx_path = run_dir / f"executive_summary_{org_domain}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    customer_name = str(summary.get("customer_name", "unknown customer"))
    tenant_key = str(summary.get("tenant_key", "unknown tenant"))
    assessment_profile = str(summary.get("assessment_profile", "unknown"))
    org_domain = _get_org_domain(summary)
    run_stamp = run_dir.name

    devices = _get_check_summary(summary, "devices")
    alerts = _get_check_summary(summary, "alerts")
    policy_posture = _get_check_summary(summary, "policy_posture")
    watchlists = _get_check_summary(summary, "watchlists")
    watchlist_effectiveness = _get_check_summary(summary, "watchlist_effectiveness")
    drift = _get_check_summary(summary, "policy_drift")
    policy_settings = _get_check_summary(summary, "policy_settings")
    core_prevention = _get_check_summary(summary, "core_prevention_settings")
    alert_workflow = _get_check_summary(summary, "alert_workflow")
    endpoint_status = _get_check_summary(summary, "endpoint_status")
    api_connector_use = _get_check_summary(summary, "api_connector_use")
    banned_hashes = _get_check_summary(summary, "banned_hashes_age")
    daily_alerts = _get_check_summary(summary, "daily_alerts_threat_scores")
    sensor_quality = _get_check_summary(summary, "sensor_coverage_quality")
    alert_quality = _get_check_summary(summary, "alert_quality")
    user_logins = _get_check_summary(summary, "user_logins")
    policy_efficacy = _get_check_summary(summary, "policy_efficacy")
    policy_tuning = _get_check_summary(summary, "policy_tuning_analysis")
    live_query = _get_check_summary(summary, "live_query_audit_remediation")
    pending_appendix_links: list[tuple[Any, str]] = []

    appendix_started = False
    appendix_tables: list[tuple[list[list[Any]], list[str], str, int]] = []

    def _start_appendix() -> None:
        nonlocal appendix_started
        if appendix_started:
            return
        _add_appendix_divider_slide(prs)
        appendix_started = True

    def _queue_appendix_table(
        all_rows: list[list[Any]],
        headers: list[str],
        title: str,
        rows_per_slide: int | None = None,
    ) -> None:
        if all_rows:
            effective_rows_per_slide = rows_per_slide
            if effective_rows_per_slide is None:
                effective_rows_per_slide = _appendix_rows_per_slide(headers)
            appendix_tables.append((all_rows, headers, title, effective_rows_per_slide))

    _add_standard_cover_slide(
        prs,
        deck_title="Carbon Black Cloud\nHealth Check - Executive Deck",
        customer_name=customer_name,
        tenant_key=tenant_key,
        assessment_profile=assessment_profile,
        run_stamp=run_stamp,
        score_text=_format_score(summary),
        summary_text="High-level posture, trend, and remediation priorities for leadership review.",
        right_cards=[
            {
                "title": "Org Domain",
                "value": org_domain or "n/a",
                "detail": f"Tenant ID {summary.get('tenant_id', 'n/a')}",
                "accent": "22C55E",
                "value_font_size_max": 16,
            },
            {
                "title": "Health Status",
                "value": _format_score(summary),
                "detail": f"Errors {len(summary.get('errors', []))}  |  Warnings {len(summary.get('warnings', []))}",
                "accent": "F59E0B",
                "value_font_size_max": 16,
            },
            {
                "title": "Key Signals",
                "value": f"{_format_count(alerts.get('total_alerts_30d'))} alerts",
                "detail": f"{_format_count(devices.get('total_devices'))} devices",
                "accent": "60A5FA",
            },
            {
                "title": "Entitlements",
                "value": _format_count(len(summary.get("products", {}).get("final", []))),
                "detail": _summarize_entitlements(summary.get("products", {}).get("final", []), max_items=None),
                "accent": "8B5CF6",
                "detail_word_wrap": True,
            },
        ],
    )

    _add_score_breakdown_slide(prs, summary, assessment_profile)
    _add_metric_definitions_slide(
        prs,
        sensor_quality=sensor_quality,
        alert_quality=alert_quality,
        policy_efficacy=policy_efficacy,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F6F8FB")
    _add_section_header(slide, "Executive Snapshot", "Visual baseline: signal mix, endpoint posture, and alert trend in one frame.")
    
    alerts_api_total = int(alerts.get("total_alerts_in_api", 0))
    alerts_processed = int(alerts.get("total_alerts_processed", alerts.get("total_alerts_30d", 0)))
    alerts_card_detail = f"High severity {_format_count(alerts.get('high_severity_alerts'))}"
    if alerts_api_total > alerts_processed and alerts_processed > 0:
        alerts_card_detail = f"API total {_format_count(alerts_api_total)}"
    
    metrics = [
        ("Devices", _format_count(devices.get("total_devices")), f"Active 7d {_format_ratio(devices.get('active_ratio_last_7d'))}", "2563EB"),
        ("Alerts", _format_count(alerts_processed), alerts_card_detail, "DC2626"),
        ("Policies", _format_count(policy_posture.get("total_policies")), f"Blocking rules {_format_ratio(policy_posture.get('blocking_rule_ratio'))}", "0F766E"),
        (
            "Watchlists",
            _format_count(watchlists.get("total_watchlists")),
            f"Alerting {_format_count(watchlists.get('alerting_enabled_watchlists'))} | Reports {_format_count(watchlists.get('total_watchlist_reports'))}",
            "7C3AED",
        ),
        ("Policy Drift", _format_count(drift.get("changed_count")), f"Detected {drift.get('drift_detected', False)}", "F59E0B"),
        ("Health Score", str(summary.get("health", {}).get("score", "n/a")), str(summary.get("health", {}).get("status", "n/a")), "111827"),
    ]

    positions = [
        (0.65, 1.35), (3.25, 1.35), (5.85, 1.35),
        (0.65, 2.95), (3.25, 2.95), (5.85, 2.95),
    ]
    for (title, value, detail, accent), (left, top) in zip(metrics, positions):
        _add_card(slide, left, top, 2.4, 1.25, title, value, detail, accent=accent)

    endpoint_counts = endpoint_status.get("status_counts", {}) if isinstance(endpoint_status, dict) else {}
    endpoint_items = _sorted_mapping_items(endpoint_counts, limit=6, sort_by_value=True)
    _add_doughnut_chart(
        slide,
        8.45,
        1.2,
        4.25,
        2.25,
        "Endpoint State Mix",
        [name for name, _ in endpoint_items],
        [value for _, value in endpoint_items],
        series_name="Endpoints",
    )

    severity_items = _sorted_mapping_items(alerts.get("severity_breakdown", {}), limit=6, sort_by_value=True)
    _add_column_chart(
        slide,
        8.45,
        3.55,
        4.25,
        2.1,
        "Top Alert Severities",
        [name for name, _ in severity_items],
        [value for _, value in severity_items],
        series_name="Alerts",
    )

    daily_counts = daily_alerts.get("daily_alert_counts", {}) if isinstance(daily_alerts, dict) else {}
    trend_labels, trend_values = _daily_alert_trend_data(daily_counts, max_points=35)
    trend_title = "Daily Alerts Trend"
    if alerts_api_total > alerts_processed and alerts_processed > 0:
        trend_title = f"Daily Alerts Trend (sampled {_format_count(alerts_processed)} of {_format_count(alerts_api_total)})"
    _add_line_chart(slide, 0.65, 4.35, 7.65, 2.05, trend_title, trend_labels, trend_values, series_name="Alerts/day")

    live_query_check = summary.get("checks", {}).get("live_query_audit_remediation", {})
    live_query_status = str(live_query_check.get("status", "")).strip().lower() if isinstance(live_query_check, dict) else ""
    if live_query_status == "ok":
        _add_live_query_audit_slide(prs, live_query)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "Policy and Prevention Posture", "Where protection is enforced, where it is only reporting, and where the policy stack is still too permissive.")

    _add_card(slide, 0.65, 1.3, 2.8, 1.3, "Enabled Policies", _format_count(policy_posture.get("enabled_policies")), f"of {_format_count(policy_posture.get('total_policies'))}", accent="0F766E")
    _add_card(slide, 3.6, 1.3, 2.8, 1.3, "Enabled Rules", _format_count(policy_posture.get("enabled_rules")), f"of {_format_count(policy_posture.get('total_rules'))}", accent="2563EB")
    _add_card(slide, 6.55, 1.3, 2.8, 1.3, "Blocking Rules", _format_count(policy_posture.get("blocking_rules")), f"ratio {_format_ratio(policy_posture.get('blocking_rule_ratio'))}", accent="DC2626")
    _add_card(slide, 9.5, 1.3, 3.15, 1.3, "Core Prevention", _format_count(len(core_prevention.get("alert_only_policies", []))), "alert-only policies", accent="F59E0B")

    if _org_defense_rules_disabled(summary):
        _add_eedr_not_applicable(slide, 0.75, 2.85, 5.9, 2.15, "Rule actions")
        _add_eedr_not_applicable(slide, 6.95, 2.85, 2.85, 2.15, "Core prevention modes")
    else:
        rule_action_items = _sorted_mapping_items(policy_posture.get("rule_action_breakdown", {}), limit=6, sort_by_value=True)
        _add_column_chart(
            slide,
            0.75,
            2.85,
            5.9,
            2.15,
            "Rule Actions",
            [name for name, _ in rule_action_items],
            [value for _, value in rule_action_items],
            series_name="Rules",
        )

        category_modes = core_prevention.get("category_mode_breakdown", {}) if isinstance(core_prevention, dict) else {}
        core_mode = category_modes.get("core_prevention", {}) if isinstance(category_modes, dict) else {}
        mode_items = _sorted_mapping_items(core_mode, limit=3, sort_by_value=False)
        _add_doughnut_chart(
            slide,
            6.95,
            2.85,
            2.85,
            2.15,
            "Core Prevention Modes",
            [name for name, _ in mode_items],
            [value for _, value in mode_items],
            series_name="Rules",
        )

    _add_textbox(
        slide,
        0.75,
        6.15,
        11.9,
        0.22,
        "Policy Settings of Concern moved to next slide for full labels and readability.",
        font_size=10,
        color=_hex_color("55606E"),
    )

    setting_items = _policy_settings_concern_items(policy_settings)
    setting_rows = [[_policy_setting_display_name(name), int(_to_numeric(value))] for name, value in setting_items]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(
        slide,
        "Policy Settings of Concern",
        "Counts show how many policies match each concern condition. The text explains why each setting matters.",
    )

    _add_textbox(
        slide,
        0.65,
        1.15,
        12.0,
        0.25,
        "Why each setting is flagged",
        font_size=11,
        bold=True,
        color=_hex_color("475569"),
    )

    _add_bar_chart(
        slide,
        0.65,
        1.45,
        7.25,
        5.9,
        "Concern Counts by Setting",
        [row[0] for row in setting_rows],
        [row[1] for row in setting_rows],
        series_name="Policies",
    )

    rationale_items = [
        (
            f"{_policy_setting_display_name(name)} ({_format_count(value)})",
            _policy_setting_rationale(name),
        )
        for name, value in reversed(setting_items)
    ]
    _add_bold_label_bullet_list(
        slide,
        8.15,
        1.5,
        4.45,
        5.65,
        rationale_items,
        font_size=10,
        color=_hex_color("23313F"),
    )

    policy_tuning_check = summary.get("checks", {}).get("policy_tuning_analysis", {})
    policy_tuning_status = str(policy_tuning_check.get("status", "")).strip().lower() if isinstance(policy_tuning_check, dict) else ""
    if policy_tuning_status == "ok":
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_background(slide, "FFFFFF")
        _add_section_header(
            slide,
            "Policy Analysis (Good/Better/Best)",
            "High-level maturity view. Detailed policy actions are in the appendix.",
        )

        current_tier = str(policy_tuning.get("current_tier", "unknown")).upper()
        score = policy_tuning.get("score_0_100", "n/a")
        review_cadence = str(policy_tuning.get("review_cadence", "quarterly_or_biannual")).replace("_", " ")
        policy_metrics = policy_tuning.get("metrics", {}) if isinstance(policy_tuning, dict) else {}

        _add_card(slide, 0.65, 1.25, 2.9, 1.15, "Maturity Tier", current_tier, "framework status", accent="0F766E")
        _add_card(slide, 3.8, 1.25, 2.9, 1.15, "Policy Score", str(score), "0-100", accent="2563EB")
        _add_card(
            slide,
            6.95,
            1.25,
            2.9,
            1.15,
            "Gaps",
            _format_count(policy_metrics.get("policies_with_enforcement_gaps", 0)),
            "policies",
            accent="DC2626",
        )
        _add_card(slide, 10.1, 1.25, 2.55, 1.15, "Review Cadence", review_cadence, "recommended", accent="7C3AED", value_font_size_max=14)

        gates = policy_tuning.get("gates", {}) if isinstance(policy_tuning, dict) else {}
        gate_rows: list[list[Any]] = []
        for gate_name in ["good", "better", "best"]:
            gate = gates.get(gate_name, {}) if isinstance(gates, dict) else {}
            passed = bool(gate.get("pass", False)) if isinstance(gate, dict) else False
            criteria = str(gate.get("criteria", "n/a")) if isinstance(gate, dict) else "n/a"
            gate_rows.append([gate_name.upper(), "MET" if passed else "NOT MET", criteria])

        _add_table(
            slide,
            0.65,
            2.65,
            12.0,
            1.9,
            ["Tier Gate", "Gate Status", "Criteria"],
            gate_rows or [["n/a", "n/a", "Policy tuning analysis unavailable."]],
        )

        top_gaps = policy_tuning.get("top_gaps", []) if isinstance(policy_tuning, dict) else []
        next_actions = policy_tuning.get("next_actions", []) if isinstance(policy_tuning, dict) else []
        gap_rows = [["Gap", str(item)] for item in (top_gaps[:3] if isinstance(top_gaps, list) else [])]
        gap_rows.extend([["Action", str(item)] for item in (next_actions[:3] if isinstance(next_actions, list) else [])])
        _add_table(
            slide,
            0.65,
            4.75,
            12.0,
            2.0,
            ["Type", "Detail"],
            gap_rows or [["Info", "No policy gaps/actions were generated."]],
        )

        policy_rec_rows_exec = _recommendation_rows(summary.get("recommendations", []), policy_only=True)
        if policy_rec_rows_exec:
            policy_analysis_footer = _add_textbox(
                slide,
                0.65,
                7.05,
                9.8,
                0.18,
                f"{len(policy_rec_rows_exec)} detailed policy recommendations are available in the appendix. Click to open.",
                font_size=8,
                color=_hex_color("2563EB"),
                underline=True,
            )
            pending_appendix_links.append((policy_analysis_footer, "Policy Recommendations (Page 1/"))

        policy_plan_items_exec = _policy_maturity_action_plan_items(policy_tuning, attention_only=True)
        if policy_plan_items_exec:
            _add_textbox(
                slide,
                0.65,
                6.84,
                11.2,
                0.18,
                f"{len(policy_plan_items_exec)} policies need maturity attention. Per-policy action-plan appendix is available only in the technical deck.",
                font_size=8,
                color=_hex_color("55606E"),
            )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(slide, "Operational Risk Areas", "These are the places where the environment drifts from healthy to annoying, noisy, or slow to respond.")

    _add_card(slide, 0.65, 1.3, 2.9, 1.25, "Alert Workflow", _format_count(alert_workflow.get("closed_alerts")), f"Open >72h {_format_count(alert_workflow.get('open_alerts_over_72h'))}", accent="7C3AED")
    _add_card(slide, 3.8, 1.3, 2.9, 1.25, "Endpoint Status", _format_count(endpoint_status.get("non_active_total")), "non-active devices", accent="EF4444")
    _add_card(slide, 6.95, 1.3, 2.9, 1.25, "Sensor Quality", _format_count(sensor_quality.get("stale_sensors_over_7d")), "stale sensors", accent="F59E0B")
    _add_card(slide, 10.1, 1.3, 2.55, 1.25, "API Connectors", _format_count(api_connector_use.get("active_connector_count")), f"Dormant {_format_count(api_connector_use.get('dormant_integration_count'))}", accent="0F766E")

    _add_card(slide, 0.65, 2.95, 2.9, 1.25, "Banned Hashes", _format_count(banned_hashes.get("older_than_365d")), "older than 365d", accent="B45309")
    _add_card(slide, 3.8, 2.95, 2.9, 1.25, "User Logins", _format_count(user_logins.get("users_with_login_events")), f"Dormant 30d {_format_count(user_logins.get('dormant_over_30d_count'))}", accent="2563EB")
    _add_card(slide, 6.95, 2.95, 2.9, 1.25, "Alert Quality", _format_count(alert_quality.get("repeated_detection_keys_over_3")), f"Noise {_format_ratio(alert_quality.get('noise_ratio'))}", accent="DC2626")
    _add_card(slide, 10.1, 2.95, 2.55, 1.25, "Efficacy", _format_count(policy_efficacy.get("total_policy_groups")), "policy groups", accent="111827")

    workflow_status = alert_workflow.get("status_counts", {}) if isinstance(alert_workflow, dict) else {}
    workflow_items = _sorted_mapping_items(workflow_status, limit=5, sort_by_value=True)
    _add_column_chart(
        slide,
        0.65,
        4.45,
        4.0,
        2.0,
        "Alert Workflow Status",
        [name for name, _ in workflow_items],
        [value for _, value in workflow_items],
        series_name="Alerts",
    )

    sensor_items = _sorted_mapping_items(sensor_quality.get("sensor_version_breakdown", {}), limit=5, sort_by_value=True)
    _add_column_chart(
        slide,
        4.9,
        4.45,
        3.9,
        2.0,
        "Sensor Versions",
        [name for name, _ in sensor_items],
        [value for _, value in sensor_items],
        series_name="Devices",
    )

    trend_labels, trend_values = _daily_alert_trend_data(daily_alerts.get("daily_alert_counts", {}), max_points=28)
    disposition_breakdown = alert_quality.get("disposition_breakdown", {}) if isinstance(alert_quality, dict) else {}
    disposition_items = [
        ("False Positive", _to_numeric(disposition_breakdown.get("False Positive", 0))),
        ("True Positive", _to_numeric(disposition_breakdown.get("True Positive", 0))),
        ("None", _to_numeric(disposition_breakdown.get("None", 0))),
    ]
    _add_pie_chart(
        slide,
        9.05,
        4.45,
        3.65,
        2.0,
        "Alert Disposition",
        [name for name, _ in disposition_items],
        [value for _, value in disposition_items],
        series_name="Alerts",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "Prioritized Recommendations", "A compact list of the highest-value next actions. This is the slide people forward.")

    recommendations = summary.get("recommendations", [])
    recommendation_rows_all = _recommendation_rows(recommendations)
    policy_recommendation_rows = _recommendation_rows(recommendations, policy_only=True)
    non_policy_recommendation_rows = _recommendation_rows(recommendations, policy_only=False)

    policy_summary_row = _policy_summary_row(policy_tuning, len(policy_recommendation_rows))
    recommendation_rows_main = list(non_policy_recommendation_rows)
    if policy_summary_row is not None:
        recommendation_rows_main.insert(0, policy_summary_row)
    if not recommendation_rows_main:
        recommendation_rows_main = recommendation_rows_all

    rec_rows: list[list[Any]] = []
    rec_rows = recommendation_rows_main[:EXEC_RECOMMENDATIONS_ROWS_ON_SLIDE]
    
    total_recs = len([r for r in recommendations if isinstance(r, dict)]) if isinstance(recommendations, list) else 0
    shown_recs = len(rec_rows)
    
    _add_table(
        slide,
        0.55,
        1.35,
        12.2,
        4.9,
        ["Priority", "Area", "Recommendation", "Evidence"],
        rec_rows or [["n/a", "n/a", "No recommendations were generated from the summary payload.", ""]],
    )
    
    recs_footnote = _add_sampling_footnote(slide, total_recs, shown_recs, y_position=6.35)
    if total_recs > shown_recs and recs_footnote is not None:
        pending_appendix_links.append((recs_footnote, "Detailed Recommendations (Page 1/"))

    if policy_recommendation_rows:
        policy_footer = _add_textbox(
            slide,
            0.62,
            6.15,
            10.8,
            0.18,
            f"{len(policy_recommendation_rows)} policy recommendations are summarized here and detailed in the appendix. Click to open full policy recommendation list.",
            font_size=8,
            color=_hex_color("2563EB"),
            underline=True,
        )
        pending_appendix_links.append((policy_footer, "Policy Recommendations (Page 1/"))
        _queue_appendix_table(
            policy_recommendation_rows,
            ["Priority", "Area", "Recommendation", "Evidence"],
            "Policy Recommendations",
        )

    dormant_appendix = _dormant_recommendation_appendix_data(recommendations, user_logins)
    if dormant_appendix is not None:
        dormant_rows, dormant_headers, shown_dormant_users, total_dormant_users, dormant_title = dormant_appendix
        dormant_footer = _add_textbox(
            slide,
            0.62,
            7.0,
            8.5,
            0.18,
            f"Showing {shown_dormant_users} of {total_dormant_users} dormant users from Evidence. Click to open appendix full list.",
            font_size=8,
            color=_hex_color("2563EB"),
            underline=True,
        )
        pending_appendix_links.append((dormant_footer, f"{dormant_title} (Page 1/"))
        _queue_appendix_table(dormant_rows, dormant_headers, dormant_title)

    errors = summary.get("errors", [])
    warnings = summary.get("warnings", [])
    issue_rows: list[list[Any]] = []
    if isinstance(errors, list):
        for message in errors:
            issue_rows.append(["Error", message])
    if isinstance(warnings, list):
        for message in warnings:
            issue_rows.append(["Warning", message])

    issue_rows_to_show = issue_rows[:EXEC_ISSUES_ROWS_ON_SLIDE]
    _add_table(
        slide,
        0.55,
        6.35,
        12.2,
        0.95,
        ["Type", "Message"],
        issue_rows_to_show or [["Info", "No run errors or warnings were captured in summary.json."]],
    )
    if len(issue_rows) > len(issue_rows_to_show):
        issues_footer = _add_textbox(
            slide,
            0.62,
            7.18,
            7.6,
            0.18,
            f"Showing {len(issue_rows_to_show)} of {len(issue_rows)} issues. Click to open appendix full list.",
            font_size=8,
            color=_hex_color("2563EB"),
            underline=True,
        )
        pending_appendix_links.append((issues_footer, "Run Issues (Page 1/"))
        _queue_appendix_table(issue_rows, ["Type", "Message"], "Run Issues")
    
    # Add pagination slides for full recommendation list if needed
    if total_recs > shown_recs:
        full_rec_rows = recommendation_rows_all
        _queue_appendix_table(
            full_rec_rows,
            ["Priority", "Area", "Recommendation", "Evidence"],
            "Detailed Recommendations",
        )

    if appendix_tables:
        _start_appendix()
        for all_rows, headers, title, rows_per_slide in appendix_tables:
            _paginate_table_to_slides(prs, all_rows, headers, title, rows_per_slide=rows_per_slide)

    for shape, title_prefix in pending_appendix_links:
        target_slide = _find_slide_by_title_prefix(prs, title_prefix)
        if target_slide is None:
            continue
        shape.click_action.target_slide = target_slide

    prs.save(str(pptx_path))
    return pptx_path


def write_technical_pptx(run_dir: Path) -> Path:
    summary_path = run_dir / "summary.json"
    if not summary_path.exists():
        raise FileNotFoundError(f"summary.json not found in {run_dir}")

    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    org_domain = _safe_path_part(_get_org_domain(summary))
    pptx_path = run_dir / f"technical_deck_{org_domain}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    customer_name = str(summary.get("customer_name", "unknown customer"))
    tenant_key = str(summary.get("tenant_key", "unknown tenant"))
    org_domain = _get_org_domain(summary)
    assessment_profile = str(summary.get("assessment_profile", "unknown"))
    errors = summary.get("errors", [])
    warnings = summary.get("warnings", [])

    devices = _get_check_summary(summary, "devices")
    alerts = _get_check_summary(summary, "alerts")
    policy_posture = _get_check_summary(summary, "policy_posture")
    watchlists = _get_check_summary(summary, "watchlists")
    watchlist_effectiveness = _get_check_summary(summary, "watchlist_effectiveness")
    drift = _get_check_summary(summary, "policy_drift")
    policy_settings = _get_check_summary(summary, "policy_settings")
    core_prevention = _get_check_summary(summary, "core_prevention_settings")
    alert_workflow = _get_check_summary(summary, "alert_workflow")
    api_connector_use = _get_check_summary(summary, "api_connector_use")
    banned_hashes = _get_check_summary(summary, "banned_hashes_age")
    endpoint_status = _get_check_summary(summary, "endpoint_status")
    daily_alerts = _get_check_summary(summary, "daily_alerts_threat_scores")
    user_logins = _get_check_summary(summary, "user_logins")
    sensor_quality = _get_check_summary(summary, "sensor_coverage_quality")
    alert_quality = _get_check_summary(summary, "alert_quality")
    policy_efficacy = _get_check_summary(summary, "policy_efficacy")
    policy_tuning = _get_check_summary(summary, "policy_tuning_analysis")
    permissions_rule_audit = _get_check_summary(summary, "permissions_rule_audit")
    live_query = _get_check_summary(summary, "live_query_audit_remediation")
    pending_appendix_links: list[tuple[Any, str]] = []

    appendix_started = False
    appendix_tables: list[tuple[list[list[Any]], list[str], str, int]] = []
    appendix_action_plans: list[tuple[list[dict[str, Any]], str]] = []

    def _start_appendix() -> None:
        nonlocal appendix_started
        if appendix_started:
            return
        _add_appendix_divider_slide(prs)
        appendix_started = True

    def _queue_appendix_table(
        all_rows: list[list[Any]],
        headers: list[str],
        title: str,
        rows_per_slide: int | None = None,
    ) -> None:
        if all_rows:
            effective_rows_per_slide = rows_per_slide
            if effective_rows_per_slide is None:
                effective_rows_per_slide = _appendix_rows_per_slide(headers)
            appendix_tables.append((all_rows, headers, title, effective_rows_per_slide))

    def _queue_policy_maturity_action_plan(items: list[dict[str, Any]], title: str) -> None:
        if items:
            appendix_action_plans.append((items, title))

    _add_standard_cover_slide(
        prs,
        deck_title="Carbon Black Cloud\nHealth Check - Technical Deck",
        customer_name=customer_name,
        tenant_key=tenant_key,
        assessment_profile=assessment_profile,
        run_stamp=run_dir.name,
        score_text=_format_score(summary),
        summary_text="Technical breakdown of check status, time trends, policy posture, and high-signal remediation priorities.",
        right_cards=[
            {
                "title": "Org Domain",
                "value": org_domain or "n/a",
                "detail": f"Tenant ID {summary.get('tenant_id', 'n/a')}",
                "accent": "22C55E",
                "value_font_size_max": 16,
            },
            {
                "title": "Health Status",
                "value": _format_score(summary),
                "detail": f"Errors {len(errors)}  |  Warnings {len(warnings)}",
                "accent": "F59E0B",
                "value_font_size_max": 16,
            },
            {
                "title": "Checks",
                "value": _format_count(len(summary.get('checks', {}))),
                "detail": "captured sections",
                "accent": "60A5FA",
            },
            {
                "title": "Entitlements",
                "value": _format_count(len(summary.get("products", {}).get("final", []))),
                "detail": _summarize_entitlements(summary.get("products", {}).get("final", []), max_items=None),
                "accent": "8B5CF6",
                "detail_word_wrap": True,
            },
        ],
    )

    _add_score_breakdown_slide(prs, summary, assessment_profile)
    _add_metric_definitions_slide(
        prs,
        sensor_quality=sensor_quality,
        alert_quality=alert_quality,
        policy_efficacy=policy_efficacy,
    )

    check_rows = []

    for name, check in summary.get("checks", {}).items():
        if not isinstance(check, dict):
            continue
        status = str(check.get("status", "unknown")).strip().lower()
        if status not in {"ok", "error", "unavailable", "unknown", "not_applicable"}:
            status = "unknown"
        summary_fields = check.get("summary", {})

        note = str(check.get("message", "")).strip()
        if not note and isinstance(summary_fields, dict):
            note_pairs: list[str] = []
            for key, value in summary_fields.items():
                if isinstance(value, (dict, list)):
                    continue
                note_pairs.append(f"{key}={value}")
                if len(note_pairs) >= 2:
                    break
            note = ", ".join(note_pairs)

        note = note or "No additional detail in summary."
        check_rows.append([
            str(name),
            status,
            note,
        ])

    status_rank = {"error": 0, "unavailable": 1, "ok": 2, "not_applicable": 3, "unknown": 4}
    check_rows.sort(key=lambda row: (status_rank.get(str(row[1]).lower(), 3), str(row[0]).lower()))

    status_colors = {
        "ok": "22C55E",
        "unavailable": "F59E0B",
        "error": "EF4444",
        "not_applicable": "0EA5A4",
        "unknown": "64748B",
    }

    status_counts: dict[str, int] = {"error": 0, "unavailable": 0, "ok": 0, "not_applicable": 0, "unknown": 0}
    for _, status, _ in check_rows:
        status_counts[str(status)] = status_counts.get(str(status), 0) + 1
    total_checks = len(check_rows)
    reliability_ratio = (status_counts.get("ok", 0) / total_checks) if total_checks else 0.0
    run_status_display = {
        "error": "RUN ERROR",
        "unavailable": "UNAVAILABLE",
        "ok": "RUN OK",
        "not_applicable": "NOT APPLICABLE",
        "unknown": "UNKNOWN",
    }

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(
        slide,
        "Run Reliability Overview",
        "Run execution status concentration and action-required checks for this run.",
    )
    _add_card(
        slide,
        0.65,
        1.35,
        3.8,
        1.15,
        "Checks Evaluated",
        _format_count(total_checks),
        f"Run-OK ratio {reliability_ratio:.0%}",
        accent="2563EB",
    )
    _add_card(
        slide,
        4.6,
        1.35,
        3.8,
        1.15,
        "Action Required",
        _format_count(status_counts.get("error", 0) + status_counts.get("unavailable", 0)),
        f"Errors {status_counts.get('error', 0)} | Unavailable {status_counts.get('unavailable', 0)}",
        accent="DC2626",
    )
    _add_card(
        slide,
        8.55,
        1.35,
        4.1,
        1.15,
        "Run-OK Checks",
        _format_count(status_counts.get("ok", 0)),
        f"N/A {status_counts.get('not_applicable', 0)} | Unknown {status_counts.get('unknown', 0)}",
        accent="16A34A",
    )

    status_labels = ["Run Error", "Unavailable", "Run OK", "Not Applicable", "Unknown"]
    status_values = [
        float(status_counts.get("error", 0)),
        float(status_counts.get("unavailable", 0)),
        float(status_counts.get("ok", 0)),
        float(status_counts.get("not_applicable", 0)),
        float(status_counts.get("unknown", 0)),
    ]
    _add_column_chart(
        slide,
        0.65,
        2.75,
        6.2,
        3.95,
        "Status Distribution",
        status_labels,
        status_values,
        series_name="Checks",
    )

    action_rows = [
        [name, run_status_display.get(status, status.replace("_", " ").upper()), note if len(note) <= 85 else (note[:82] + "...")]
        for name, status, note in check_rows
        if status in {"error", "unavailable"}
    ][:6]
    if action_rows:
        _add_table(slide, 7.0, 2.75, 5.65, 3.95, ["Check", "Status", "Reason"], action_rows)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(
        slide,
        "Check Run Status (Prioritized)",
        "Run execution outcomes only. Checks are ordered by urgency: error, unavailable, run ok, unknown.",
    )

    _add_card(
        slide,
        0.65,
        1.35,
        3.0,
        0.95,
        "Error",
        _format_count(status_counts.get("error", 0)),
        "Immediate remediation",
        accent=status_colors["error"],
    )
    _add_card(
        slide,
        3.8,
        1.35,
        3.0,
        0.95,
        "Unavailable",
        _format_count(status_counts.get("unavailable", 0)),
        "Data/API follow-up",
        accent=status_colors["unavailable"],
    )
    _add_card(
        slide,
        6.95,
        1.35,
        2.7,
        0.95,
        "Run OK",
        _format_count(status_counts.get("ok", 0)),
        "No immediate action",
        accent=status_colors["ok"],
    )
    _add_card(
        slide,
        9.8,
        1.35,
        2.85,
        0.95,
        "Unknown",
        _format_count(status_counts.get("unknown", 0)),
        f"N/A {status_counts.get('not_applicable', 0)}",
        accent=status_colors["unknown"],
    )

    outcomes_rows_all = [
        [name, run_status_display.get(status, status.replace("_", " ").upper()), note if len(note) <= 110 else (note[:107] + "...")]
        for name, status, note in check_rows
    ]
    outcomes_rows = outcomes_rows_all[:TECH_CHECK_OUTCOMES_ROWS_ON_SLIDE]
    if outcomes_rows:
        _add_table(slide, 0.65, 2.45, 12.0, 4.6, ["Check", "Status", "Note"], outcomes_rows)
    else:
        _add_unavailable_callout(slide, 0.65, 2.45, 12.0, 4.6, "Check Outcomes", "No check rows available")

    if len(check_rows) > len(outcomes_rows):
        outcomes_footer = _add_textbox(
            slide,
            0.65,
            7.12,
            9.4,
            0.2,
            f"Showing {len(outcomes_rows)} of {len(check_rows)} checks on this slide. Click to open appendix full outcomes.",
            font_size=8,
            color=_hex_color("2563EB"),
            underline=True,
        )
        pending_appendix_links.append((outcomes_footer, "All Check Outcomes (Page 1/"))
        _queue_appendix_table(
            [
                [name, run_status_display.get(status, status.replace("_", " ").upper()), note if len(note) <= 200 else (note[:197] + "...")]
                for name, status, note in check_rows
            ],
            ["Check", "Status", "Note"],
            "All Check Outcomes",
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "Alert Volume and Mix", "Alert trend and mix stay together; endpoint platform posture is separated onto the next slide.")
    endpoint_counts = endpoint_status.get("status_counts", {}) if isinstance(endpoint_status, dict) else {}
    active_or_bypass_devices = int(
        _to_numeric(endpoint_counts.get("ACTIVE", 0))
        + _to_numeric(endpoint_counts.get("BYPASS", 0))
    )
    if active_or_bypass_devices <= 0:
        total_devices = _to_numeric(devices.get("total_devices", 0))
        deregistered_devices = _to_numeric(endpoint_counts.get("DEREGISTERED", 0))
        active_or_bypass_devices = max(0, int(total_devices - deregistered_devices))
    _add_card(
        slide,
        0.65,
        1.2,
        2.75,
        1.2,
        "Devices",
        _format_count(active_or_bypass_devices),
        f"Excludes deregistered | Active 7d {_format_ratio(devices.get('active_ratio_last_7d'))}",
        accent="2563EB",
    )
    
    _tech_alerts_api_total = int(alerts.get("total_alerts_in_api", 0))
    _tech_alerts_processed = int(alerts.get("total_alerts_processed", alerts.get("total_alerts_30d", 0)))
    _tech_alerts_detail = f"High severity {_format_count(alerts.get('high_severity_alerts'))}"
    if _tech_alerts_api_total > _tech_alerts_processed and _tech_alerts_processed > 0:
        _tech_alerts_detail = f"API total {_format_count(_tech_alerts_api_total)}"
    _add_card(slide, 3.55, 1.2, 2.75, 1.2, "Alerts", _format_count(_tech_alerts_processed), _tech_alerts_detail, accent="DC2626")
    
    _add_card(slide, 6.45, 1.2, 2.75, 1.2, "Avg Daily Alerts", _format_count(daily_alerts.get("avg_daily_alerts")), f"Spike days {_format_count(len(daily_alerts.get('spike_days', [])))}", accent="7C3AED")
    _add_card(slide, 9.35, 1.2, 3.3, 1.2, "Endpoint Status", _format_count(endpoint_status.get("non_active_total")), "non-active devices", accent="0F766E")

    trend_labels, trend_values = _daily_alert_trend_data(daily_alerts.get("daily_alert_counts", {}), max_points=45)
    _tech_trend_title = "Daily Alerts Over Time"
    if _tech_alerts_api_total > _tech_alerts_processed and _tech_alerts_processed > 0:
        _tech_trend_title = f"Daily Alerts (sampled {_format_count(_tech_alerts_processed)} of {_format_count(_tech_alerts_api_total)})"
    _add_line_chart(slide, 0.65, 2.65, 7.25, 3.8, _tech_trend_title, trend_labels, trend_values, series_name="Alerts/day")

    severity_items = _sorted_mapping_items(alerts.get("severity_breakdown", {}), limit=6, sort_by_value=True)
    _add_column_chart(
        slide,
        8.15,
        2.65,
        4.55,
        1.85,
        "Top Severities",
        [name for name, _ in severity_items],
        [value for _, value in severity_items],
        series_name="Alerts",
        show_data_labels=False,
    )

    alert_type_items = _sorted_mapping_items(alerts.get("type_breakdown", {}), limit=5, sort_by_value=True)
    _add_column_chart(
        slide,
        8.15,
        4.6,
        4.55,
        1.85,
        "Alert Type Breakdown",
        [name for name, _ in alert_type_items],
        [value for _, value in alert_type_items],
        series_name="Alerts",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(slide, "Endpoint Platform and Sensor Health", "Device OS mix and sensor-version spread live together so this reads as fleet posture, not alert context.")

    _add_card(slide, 0.65, 1.2, 2.75, 1.2, "Devices", _format_count(devices.get("total_devices")), f"Active 7d {_format_ratio(devices.get('active_ratio_last_7d'))}", accent="2563EB")
    _add_card(slide, 3.55, 1.2, 2.75, 1.2, "Active 7d", _format_count(devices.get("active_last_7d")), "reporting recently", accent="0F766E")
    _add_card(slide, 6.45, 1.2, 2.75, 1.2, "Stale Sensors", _format_count(sensor_quality.get("stale_sensors_over_7d")), "over 7 days", accent="F59E0B")
    _add_card(slide, 9.35, 1.2, 3.3, 1.2, "Non-active Endpoints", _format_count(endpoint_status.get("non_active_total")), "status exceptions", accent="DC2626")

    os_items = _sorted_mapping_items(devices.get("os_breakdown", {}), limit=6, sort_by_value=True)
    _add_doughnut_chart(
        slide,
        0.65,
        2.85,
        5.95,
        3.35,
        "Device OS Mix",
        [name for name, _ in os_items],
        [value for _, value in os_items],
        series_name="Devices",
    )

    sensor_items = _sorted_mapping_items(sensor_quality.get("sensor_version_breakdown", {}), limit=10, sort_by_value=True)
    _add_bar_chart(
        slide,
        6.75,
        2.85,
        5.95,
        3.35,
        "Top Sensor Versions",
        [_truncate_chart_label(name, max_length=26) for name, _ in sensor_items],
        [value for _, value in sensor_items],
        series_name="Devices",
        show_data_labels=False,
        category_font_size=8,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "Alert Top 10 Breakouts", "This detail belongs in the technical deck: ATT&CK techniques, alerted processes, and reputation mix.")

    attack_technique_items = _sorted_mapping_items(alerts.get("attack_technique_breakdown", {}), limit=10, sort_by_value=True)
    _add_bar_chart(
        slide,
        0.65,
        1.35,
        12.05,
        2.05,
        "Top 10 Alerted ATT&CK Techniques",
        [_truncate_chart_label(name, max_length=48) for name, _ in attack_technique_items],
        [value for _, value in attack_technique_items],
        series_name="Alerts",
        show_data_labels=False,
        category_font_size=8,
    )

    process_items = _sorted_mapping_items(alerts.get("process_breakdown", {}), limit=10, sort_by_value=True)
    _add_bar_chart(
        slide,
        0.65,
        3.9,
        5.95,
        2.35,
        "Top 10 Alerted Processes",
        [_truncate_chart_label(name, max_length=28) for name, _ in process_items],
        [value for _, value in process_items],
        series_name="Alerts",
        show_data_labels=False,
        category_font_size=8,
    )

    reputation_items = _sorted_mapping_items(alerts.get("reputation_breakdown", {}), limit=10, sort_by_value=True)
    _add_bar_chart(
        slide,
        6.75,
        3.9,
        5.95,
        2.35,
        "Top 10 Alert Reputations",
        [_truncate_chart_label(name, max_length=24) for name, _ in reputation_items],
        [value for _, value in reputation_items],
        series_name="Alerts",
        show_data_labels=False,
        category_font_size=8,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F6F8FB")
    _add_section_header(slide, "Policy Posture and Drift", "Policy volume, rule actions, drift, and settings that matter to the control plane.")
    _add_card(slide, 0.65, 1.2, 2.75, 1.2, "Policies", _format_count(policy_posture.get("total_policies")), f"Enabled {_format_count(policy_posture.get('enabled_policies'))}", accent="0F766E")
    _add_card(slide, 3.55, 1.2, 2.75, 1.2, "Rules", _format_count(policy_posture.get("total_rules")), f"Blocking {_format_count(policy_posture.get('blocking_rules'))}", accent="2563EB")
    _add_card(slide, 6.45, 1.2, 2.75, 1.2, "Drifted Policies", _format_count(drift.get("changed_count")), f"Detected {drift.get('drift_detected', False)}", accent="F59E0B")
    _add_card(slide, 9.35, 1.2, 3.3, 1.2, "Alert-only Core Prevention", _format_count(len(core_prevention.get('alert_only_policies', []))), "policies", accent="7C3AED")
    if _org_defense_rules_disabled(summary):
        _add_eedr_not_applicable(slide, 0.65, 2.85, 4.05, 2.15, "Rule action counts")
    else:
        rule_action_items = _sorted_mapping_items(policy_posture.get("rule_action_breakdown", {}), limit=6, sort_by_value=True)
        _add_column_chart(
            slide,
            0.65,
            2.85,
            4.05,
            2.15,
            "Rule Action Counts",
            [name for name, _ in rule_action_items],
            [value for _, value in rule_action_items],
            series_name="Rules",
        )

    if _org_defense_rules_disabled(summary):
        _add_eedr_not_applicable(slide, 4.95, 2.85, 2.65, 2.15, "Core prevention modes")
    else:
        core_modes = core_prevention.get("category_mode_breakdown", {}).get("core_prevention", {}) if isinstance(core_prevention, dict) else {}
        core_mode_items = _sorted_mapping_items(core_modes, limit=3, sort_by_value=False)
        _add_doughnut_chart(
            slide,
            4.95,
            2.85,
            2.65,
            2.15,
            "Core Prevention Modes",
            [name for name, _ in core_mode_items],
            [value for _, value in core_mode_items],
            series_name="Rules",
        )

    setting_items = _policy_settings_concern_items(policy_settings)

    drift_details = drift.get("changed_policy_details", []) if isinstance(drift, dict) else []
    drift_rows = _rows_from_list(drift_details, ["policy_name", "change_count"], 6)
    _add_table(slide, 0.65, 5.2, 12.05, 1.3, ["Drifted Policy", "Changed Fields"], drift_rows or [["No drifted policies", 0]])
    if isinstance(drift_details, list) and len(drift_details) > len(drift_rows):
        full_drift_rows = _rows_from_list(drift_details, ["policy_name", "change_count"])
        _queue_appendix_table(full_drift_rows, ["Drifted Policy", "Changed Fields"], "Policy Drift Details")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(
        slide,
        "Policy Settings of Concern",
        "Counts show how many policies match each concern condition. The text explains why each setting matters.",
    )
    _add_bar_chart(
        slide,
        0.65,
        1.45,
        7.25,
        5.9,
        "Concern Counts by Setting",
        [_policy_setting_display_name(name) for name, _ in setting_items],
        [value for _, value in setting_items],
        series_name="Policies",
    )

    rationale_items = [
        (
            f"{_policy_setting_display_name(name)} ({_format_count(value)})",
            _policy_setting_rationale(name),
        )
        for name, value in reversed(setting_items)
    ]
    _add_bold_label_bullet_list(
        slide,
        8.15,
        1.5,
        4.45,
        5.65,
        rationale_items,
        font_size=10,
        color=_hex_color("23313F"),
    )

    policy_tuning_check = summary.get("checks", {}).get("policy_tuning_analysis", {})
    policy_tuning_status = str(policy_tuning_check.get("status", "")).strip().lower() if isinstance(policy_tuning_check, dict) else ""
    if policy_tuning_status == "ok":
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_background(slide, "FFFFFF")
        _add_section_header(
            slide,
            "Policy Analysis (Good/Better/Best)",
            "High-level maturity checkpoint in the policy section. Full policy actions live in the appendix.",
        )

        current_tier = str(policy_tuning.get("current_tier", "unknown")).upper()
        score = policy_tuning.get("score_0_100", "n/a")
        review_cadence = str(policy_tuning.get("review_cadence", "quarterly_or_biannual")).replace("_", " ")
        policy_metrics = policy_tuning.get("metrics", {}) if isinstance(policy_tuning, dict) else {}

        _add_card(slide, 0.65, 1.2, 2.75, 1.2, "Maturity Tier", current_tier, "framework status", accent="0F766E")
        _add_card(slide, 3.55, 1.2, 2.75, 1.2, "Policy Score", str(score), "0-100", accent="2563EB")
        _add_card(
            slide,
            6.45,
            1.2,
            2.75,
            1.2,
            "Gaps",
            _format_count(policy_metrics.get("policies_with_enforcement_gaps", 0)),
            "policies",
            accent="DC2626",
        )
        _add_card(slide, 9.35, 1.2, 3.3, 1.2, "Review Cadence", review_cadence, "recommended", accent="7C3AED", value_font_size_max=14)

        gates = policy_tuning.get("gates", {}) if isinstance(policy_tuning, dict) else {}
        gate_rows: list[list[Any]] = []
        for gate_name in ["good", "better", "best"]:
            gate = gates.get(gate_name, {}) if isinstance(gates, dict) else {}
            passed = bool(gate.get("pass", False)) if isinstance(gate, dict) else False
            criteria = str(gate.get("criteria", "n/a")) if isinstance(gate, dict) else "n/a"
            gate_rows.append([gate_name.upper(), "MET" if passed else "NOT MET", criteria])

        _add_table(
            slide,
            0.65,
            2.75,
            12.0,
            1.9,
            ["Tier Gate", "Gate Status", "Criteria"],
            gate_rows or [["n/a", "n/a", "Policy tuning analysis unavailable."]],
        )

        top_gaps = policy_tuning.get("top_gaps", []) if isinstance(policy_tuning, dict) else []
        next_actions = policy_tuning.get("next_actions", []) if isinstance(policy_tuning, dict) else []
        gap_rows = [["Gap", str(item)] for item in (top_gaps[:3] if isinstance(top_gaps, list) else [])]
        gap_rows.extend([["Action", str(item)] for item in (next_actions[:3] if isinstance(next_actions, list) else [])])
        _add_table(
            slide,
            0.65,
            4.85,
            12.0,
            1.9,
            ["Type", "Detail"],
            gap_rows or [["Info", "No policy gaps/actions were generated."]],
        )

        policy_rec_rows_tech_summary = _recommendation_rows(summary.get("recommendations", []), policy_only=True)
        if policy_rec_rows_tech_summary:
            policy_analysis_footer_tech = _add_textbox(
                slide,
                0.65,
                7.05,
                9.8,
                0.18,
                f"{len(policy_rec_rows_tech_summary)} detailed policy recommendations are available in the appendix. Click to open.",
                font_size=8,
                color=_hex_color("2563EB"),
                underline=True,
            )
            pending_appendix_links.append((policy_analysis_footer_tech, "Policy Recommendations (Page 1/"))

        policy_plan_items_tech = _policy_maturity_action_plan_items(policy_tuning, attention_only=True)
        if policy_plan_items_tech:
            policy_plan_footer_tech = _add_textbox(
                slide,
                0.65,
                6.84,
                10.8,
                0.18,
                f"{len(policy_plan_items_tech)} policies need maturity attention. Click for per-policy next-level action plan.",
                font_size=8,
                color=_hex_color("2563EB"),
                underline=True,
            )
            pending_appendix_links.append((policy_plan_footer_tech, "Policy Maturity Action Plan"))
            _queue_policy_maturity_action_plan(policy_plan_items_tech, "Policy Maturity Action Plan")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "Operational Hygiene", "Workflow, sensor, alert noise, user login, and connector hygiene signals.")
    _add_card(slide, 0.65, 1.2, 2.75, 1.2, "Closed Alerts", _format_count(alert_workflow.get("closed_alerts")), f"Median close {_format_count(alert_workflow.get('median_time_to_close_hours'))}h", accent="22C55E")
    _add_card(slide, 3.55, 1.2, 2.75, 1.2, "Open >72h", _format_count(alert_workflow.get("open_alerts_over_72h")), "workflow backlog", accent="EF4444")
    _add_card(slide, 6.45, 1.2, 2.75, 1.2, "Stale Sensors", _format_count(sensor_quality.get("stale_sensors_over_7d")), "over 7 days", accent="F59E0B")
    _add_card(slide, 9.35, 1.2, 3.3, 1.2, "Alert Noise", _format_ratio(alert_quality.get("noise_ratio")), f"Repeated alert signatures { _format_count(alert_quality.get('repeated_detection_keys_over_3')) }", accent="7C3AED")
    workflow_items = _sorted_mapping_items(alert_workflow.get("status_counts", {}), limit=6, sort_by_value=True)
    _add_column_chart(
        slide,
        0.65,
        2.85,
        3.95,
        2.05,
        "Workflow Status",
        [name for name, _ in workflow_items],
        [value for _, value in workflow_items],
        series_name="Alerts",
    )

    disposition_breakdown = alert_quality.get("disposition_breakdown", {}) if isinstance(alert_quality, dict) else {}
    disposition_items = [
        ("False Positive", _to_numeric(disposition_breakdown.get("False Positive", 0))),
        ("True Positive", _to_numeric(disposition_breakdown.get("True Positive", 0))),
        ("None", _to_numeric(disposition_breakdown.get("None", 0))),
    ]
    _add_pie_chart(
        slide,
        4.85,
        2.85,
        3.95,
        2.05,
        "Alert Disposition",
        [name for name, _ in disposition_items],
        [value for _, value in disposition_items],
        series_name="Alerts",
    )

    dormant_counts = {
        "7d+": _to_numeric(user_logins.get("dormant_over_7d_count", 0)) if isinstance(user_logins, dict) else 0.0,
        "30d+": _to_numeric(user_logins.get("dormant_over_30d_count", 0)) if isinstance(user_logins, dict) else 0.0,
        "60d+": _to_numeric(user_logins.get("dormant_over_60d_count", 0)) if isinstance(user_logins, dict) else 0.0,
    }
    dormant_items = list(dormant_counts.items())
    _add_column_chart(
        slide,
        0.65,
        5.05,
        5.95,
        1.45,
        "Dormant Users by Age",
        [name for name, _ in dormant_items],
        [value for _, value in dormant_items],
        series_name="Users",
    )

    unresolved_buckets = alert_quality.get("unresolved_aging_buckets", {}) if isinstance(alert_quality, dict) else {}
    if isinstance(unresolved_buckets, dict):
        age_labels = ["0-6d", "7-29d", "30+d"]
        age_items = [(label, _to_numeric(unresolved_buckets.get(label, 0))) for label in age_labels]
        if all(value == 0 for _, value in age_items):
            age_items = _sorted_mapping_items(unresolved_buckets, limit=6, sort_by_value=False)
    else:
        age_items = []
    if not age_items:
        age_items = [("No data", 0.0)]

    _add_column_chart(
        slide,
        6.75,
        5.05,
        5.95,
        1.45,
        "Unresolved Alert Aging",
        [name for name, _ in age_items],
        [value for _, value in age_items],
        series_name="Alerts",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "API Hygiene", "Connector activity, source IP ownership, and origin country for API access sessions.")

    connector_rows = _rows_from_list(api_connector_use.get("active_connectors", []), ["connector_id", "api_access_level_type", "session_count", "ip_addresses"], 8)
    connector_counts = {str(row[0]): _to_numeric(row[2]) for row in connector_rows if len(row) >= 3}
    connector_items = _sorted_mapping_items(connector_counts, limit=5, sort_by_value=True)

    top_source_ips = api_connector_use.get("top_source_ips", []) if isinstance(api_connector_use, dict) else []
    source_ip_rows = _rows_from_list(top_source_ips, ["ip", "org", "country", "session_count"], 8)
    source_ip_counts = {str(row[0]): _to_numeric(row[3]) for row in source_ip_rows if len(row) == 4}
    source_country_counts: dict[str, float] = {}
    for row in source_ip_rows:
        if len(row) != 4:
            continue
        country = str(row[2] or "unknown")
        source_country_counts[country] = source_country_counts.get(country, 0.0) + _to_numeric(row[3])

    _add_card(slide, 0.65, 1.2, 2.75, 1.2, "Connector Sessions", _format_count(api_connector_use.get("connector_session_events")), f"Active {_format_count(api_connector_use.get('active_connector_count'))}", accent="2563EB")
    _add_card(slide, 3.55, 1.2, 2.75, 1.2, "Unique Source IPs", _format_count(api_connector_use.get("unique_source_ips")), "observed in audit logs", accent="0F766E")
    _add_card(slide, 6.45, 1.2, 2.75, 1.2, "Top IP Count", _format_count(len(source_ip_rows)), "shown in deck", accent="7C3AED")
    _add_card(slide, 9.35, 1.2, 3.3, 1.2, "Countries", _format_count(len(source_country_counts)), "origin breakdown", accent="F59E0B")

    _add_column_chart(
        slide,
        0.65,
        2.85,
        3.85,
        2.1,
        "Connector Session Load",
        [name for name, _ in connector_items],
        [value for _, value in connector_items],
        series_name="Sessions",
    )

    country_items = _sorted_mapping_items(source_country_counts, limit=6, sort_by_value=True)
    _add_doughnut_chart(
        slide,
        4.7,
        2.85,
        3.1,
        2.1,
        "Session Origin Countries",
        [name for name, _ in country_items],
        [value for _, value in country_items],
        series_name="Sessions",
    )

    _add_table(
        slide,
        8.05,
        2.85,
        4.6,
        2.9,
        ["IP", "Owner", "Country", "Sessions"],
        source_ip_rows or [["n/a", "n/a", "n/a", 0]],
    )

    _add_textbox(
        slide,
        8.05,
        5.85,
        4.6,
        1.1,
        "Owner reflects the enriched organization field from IP lookup; country reflects the source IP's geolocation lookup.",
        font_size=10,
        color=_hex_color("55606E"),
    )

    _add_table(
        slide,
        0.65,
        5.25,
        7.2,
        1.65,
        ["Connector", "API Access Level Type", "Sessions", "IP Addresses"],
        connector_rows or [["n/a", "n/a", 0, "n/a"]],
    )

    live_query_check = summary.get("checks", {}).get("live_query_audit_remediation", {})
    live_query_status = str(live_query_check.get("status", "")).strip().lower() if isinstance(live_query_check, dict) else ""
    if live_query_status == "ok":
        _add_live_query_audit_slide(prs, live_query)

    watchlist_rows = _watchlist_detail_rows(watchlists.get("watchlist_details", [])) if isinstance(watchlists, dict) else []
    if watchlist_rows:
        total_watchlists = _format_count(watchlists.get("total_watchlists"))
        alerting_watchlists = _format_count(watchlists.get("alerting_enabled_watchlists"))
        enabled_without_alerting = _format_count(watchlists.get("enabled_without_alerting_watchlists"))
        total_report_count = _format_count(watchlists.get("total_watchlist_reports"))
        average_report_count = watchlists.get("average_report_count")
        if isinstance(average_report_count, (int, float)):
            average_report_detail = f"avg {average_report_count:.1f} per watchlist"
        else:
            average_report_detail = "avg n/a"

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_background(slide, "FFFFFF")
        _add_section_header(
            slide,
            "Watchlist Coverage",
            "Alerting is inferred from report_ids in watchlist metadata. Enabled watchlists with zero reports are silent.",
        )
        _add_card(slide, 0.65, 1.2, 2.75, 1.2, "Watchlists", total_watchlists, f"Enabled {_format_count(watchlists.get('enabled_watchlists'))}", accent="7C3AED")
        _add_card(slide, 3.55, 1.2, 2.75, 1.2, "Alerting Enabled", alerting_watchlists, f"of {total_watchlists}", accent="2563EB")
        _add_card(slide, 6.45, 1.2, 2.75, 1.2, "Enabled Silent", enabled_without_alerting, "enabled but no reports", accent="F59E0B")
        _add_card(slide, 9.35, 1.2, 3.3, 1.2, "Report Count", total_report_count, average_report_detail, accent="0F766E")

        _add_table(
            slide,
            0.65,
            2.85,
            12.0,
            3.7,
            ["Watchlist", "Enabled", "Alerting", "Reports"],
            watchlist_rows[:TECH_WATCHLIST_ROWS_ON_SLIDE],
        )
        watchlist_footnote = _add_sampling_footnote(
            slide,
            len(watchlist_rows),
            min(len(watchlist_rows), TECH_WATCHLIST_ROWS_ON_SLIDE),
            y_position=6.75,
        )

        if len(watchlist_rows) > TECH_WATCHLIST_ROWS_ON_SLIDE:
            if watchlist_footnote is not None:
                pending_appendix_links.append((watchlist_footnote, "Watchlist Coverage (Page 1/"))
            _queue_appendix_table(watchlist_rows, ["Watchlist", "Enabled", "Alerting", "Reports"], "Watchlist Coverage")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "F8FAFC")
    _add_section_header(slide, "Detailed Breakdown", "Top-risk lists that usually matter during engineering review.")
    risk_metric_data = {
        "Old banned hashes": _to_numeric(banned_hashes.get("older_than_365d", 0)),
        "Bypass rules": _to_numeric(permissions_rule_audit.get("bypass_rules", 0)),
        "BYPASS Critical (P1)": _to_numeric(permissions_rule_audit.get("p1_policy_count", 0)),
        "BYPASS High (P2)": _to_numeric(permissions_rule_audit.get("p2_policy_count", 0)),
    }
    risk_items = _sorted_mapping_items(risk_metric_data, limit=6, sort_by_value=False)
    _add_column_chart(
        slide,
        0.65,
        1.35,
        4.1,
        2.35,
        "Risk Metric Counts",
        [name for name, _ in risk_items],
        [value for _, value in risk_items],
        series_name="Count",
    )

    policy_breakdown = policy_efficacy.get("policy_group_breakdown", []) if isinstance(policy_efficacy, dict) else []
    if _org_defense_rules_disabled(summary):
        _add_eedr_not_applicable(slide, 5.0, 1.35, 7.7, 2.35, "Blocking-rule analysis")
    else:
        policy_blocks: dict[str, float] = {}
        for item in policy_breakdown[:12]:
            if isinstance(item, dict):
                policy_blocks[str(item.get("policy_name", "unknown"))] = _to_numeric(item.get("block_rules", 0))
        policy_items = _sorted_mapping_items(policy_blocks, limit=8, sort_by_value=True)
        _add_column_chart(
            slide,
            5.0,
            1.35,
            7.7,
            2.35,
            "Policies with Highest Blocking Rules",
            [name for name, _ in policy_items],
            [value for _, value in policy_items],
            series_name="Block rules",
            show_data_labels=False,
        )

    _add_textbox(
        slide,
        0.65,
        3.72,
        12.0,
        0.22,
        "Table purpose: per-policy enforcement mix (Monitor vs Block). Ratio = Monitor / (Monitor + Block + Unknown).",
        font_size=10,
        color=_hex_color("55606E"),
    )

    policy_breakdown = policy_efficacy.get("policy_group_breakdown", []) if isinstance(policy_efficacy, dict) else []
    full_policy_rows = _rows_from_list(policy_breakdown, ["policy_name", "monitor_rules", "block_rules", "monitor_ratio"])
    policy_rows_on_slide = _rows_from_list(
        policy_breakdown,
        ["policy_name", "monitor_rules", "block_rules", "monitor_ratio"],
        TECH_POLICY_EFFICACY_ROWS_ON_SLIDE,
    )
    total_policies = len([p for p in policy_breakdown if isinstance(p, dict)]) if isinstance(policy_breakdown, list) else 0
    shown_policies = len(policy_rows_on_slide)
    
    _add_table(slide, 0.65, 4.0, 12.0, 2.55, ["Policy", "Monitor", "Block", "Ratio"], policy_rows_on_slide)
    
    policy_footnote = _add_sampling_footnote(slide, total_policies, shown_policies, y_position=6.73)
    
    # Add pagination slides for full policy breakdown if needed
    if total_policies > shown_policies:
        if policy_footnote is not None:
            pending_appendix_links.append((policy_footnote, "Policy Efficacy Details (Page 1/"))
        _queue_appendix_table(
            full_policy_rows,
            ["Policy", "Monitor", "Block", "Ratio"],
            "Policy Efficacy Details",
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, "FFFFFF")
    _add_section_header(slide, "Recommendations", "The generated remediation list, kept in the deck for reference.")
    recommendations = summary.get("recommendations", [])
    full_rec_rows_tech = _recommendation_rows(recommendations)
    policy_recommendation_rows_tech = _recommendation_rows(recommendations, policy_only=True)
    non_policy_recommendation_rows_tech = _recommendation_rows(recommendations, policy_only=False)

    policy_summary_row_tech = _policy_summary_row(policy_tuning, len(policy_recommendation_rows_tech))
    rec_rows_pool_tech = list(non_policy_recommendation_rows_tech)
    if policy_summary_row_tech is not None:
        rec_rows_pool_tech.insert(0, policy_summary_row_tech)
    if not rec_rows_pool_tech:
        rec_rows_pool_tech = full_rec_rows_tech

    rec_rows = rec_rows_pool_tech[:TECH_RECOMMENDATIONS_ROWS_ON_SLIDE]
    
    total_recs_tech = len([r for r in recommendations if isinstance(r, dict)]) if isinstance(recommendations, list) else 0
    shown_recs_tech = len(rec_rows)
    
    _add_table(slide, 0.55, 1.25, 12.2, 5.8, ["Priority", "Area", "Recommendation", "Evidence"], rec_rows or [["n/a", "n/a", "No recommendations were generated.", ""]])
    
    recs_footnote = _add_sampling_footnote(slide, total_recs_tech, shown_recs_tech, y_position=7.18)

    if policy_recommendation_rows_tech:
        policy_footer_tech = _add_textbox(
            slide,
            0.62,
            7.0,
            10.8,
            0.18,
            f"{len(policy_recommendation_rows_tech)} policy recommendations are in the appendix. Click to open the full policy recommendation table.",
            font_size=8,
            color=_hex_color("2563EB"),
            underline=True,
        )
        pending_appendix_links.append((policy_footer_tech, "Policy Recommendations (Page 1/"))
        _queue_appendix_table(
            policy_recommendation_rows_tech,
            ["Priority", "Area", "Recommendation", "Evidence"],
            "Policy Recommendations",
        )

    dormant_appendix = _dormant_recommendation_appendix_data(recommendations, user_logins)
    if dormant_appendix is not None:
        dormant_rows, dormant_headers, shown_dormant_users, total_dormant_users, dormant_title = dormant_appendix
        dormant_footer = _add_textbox(
            slide,
            0.62,
            6.98,
            8.5,
            0.18,
            f"Showing {shown_dormant_users} of {total_dormant_users} dormant users from Evidence. Click to open appendix full list.",
            font_size=8,
            color=_hex_color("2563EB"),
            underline=True,
        )
        pending_appendix_links.append((dormant_footer, f"{dormant_title} (Page 1/"))
        _queue_appendix_table(dormant_rows, dormant_headers, dormant_title)

    # Add pagination slides for full recommendation list if needed
    if total_recs_tech > shown_recs_tech:
        if recs_footnote is not None:
            pending_appendix_links.append((recs_footnote, "Complete Recommendations (Page 1/"))
        _queue_appendix_table(
            full_rec_rows_tech,
            ["Priority", "Area", "Recommendation", "Evidence"],
            "Complete Recommendations",
        )

    if appendix_tables or appendix_action_plans:
        _start_appendix()
        for all_rows, headers, title, rows_per_slide in appendix_tables:
            _paginate_table_to_slides(prs, all_rows, headers, title, rows_per_slide=rows_per_slide)
        for items, title in appendix_action_plans:
            _paginate_policy_maturity_action_plan(prs, items, title)

    for shape, title_prefix in pending_appendix_links:
        target_slide = _find_slide_by_title_prefix(prs, title_prefix)
        if target_slide is None:
            continue
        shape.click_action.target_slide = target_slide

    prs.save(str(pptx_path))
    return pptx_path
